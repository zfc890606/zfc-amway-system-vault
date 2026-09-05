# -*- coding: utf-8 -*-
"""人生是旷野 · AI大健康-给年轻人的人生旷野课 —— 黄底黑字·明朝体+手写混搭·时尚杂志感"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 配色 / 字体 ----------
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
YELLOW_D = RGBColor(0xF0, 0xC2, 0x00)   # 水印
BLACK   = RGBColor(0x11, 0x11, 0x11)
RED     = RGBColor(0xD7, 0x00, 0x2F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x55, 0x55, 0x55)
F_SERIF = '宋体'    # 明朝体
F_KAI   = '楷体'    # 手写
F_SANS  = '黑体'    # 小字标签

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
N_TOTAL = 18

# ---------- 基础工具 ----------
def _cjk(run, font_cn):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_cn)

def _style_run(r, text, size, color, bold, font):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    _cjk(r, font)

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = YELLOW
    return s

def page_chrome(s, idx):
    # 顶部红色进度条
    pbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(W * (idx+1) / N_TOTAL), Inches(0.045))
    pbar.fill.solid(); pbar.fill.fore_color.rgb = RED
    pbar.line.fill.background()
    # 页码
    tb(s, Inches(11.7), Inches(7.02), Inches(1.3), Inches(0.4),
       f"{idx+1:02d} — {N_TOTAL:02d}", size=11, color=BLACK, font=F_SANS, align=PP_ALIGN.RIGHT)

def tb(slide, l, t, w, h, text, size=18, color=BLACK, bold=False, font=F_SERIF,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, rotation=0, ls=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = anchor
    except Exception: pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    r = p.add_run()
    _style_run(r, text, size, color, bold, font)
    if rotation: box.rotation = rotation
    return box

def tb_rich(slide, l, t, w, h, parts, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            rotation=0, ls=1.15):
    """parts: list of (text, size, color, bold, font)"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = anchor
    except Exception: pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    for (t, sz, c, b, f) in parts:
        r = p.add_run()
        _style_run(r, t, sz, c, b, f)
    if rotation: box.rotation = rotation
    return box

def para(tf, text, size=14, color=BLACK, bold=False, font=F_SERIF,
         align=PP_ALIGN.LEFT, sb=4, sa=2, ls=1.2, idx=None):
    p = tf.paragraphs[idx] if idx is not None else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = ls
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    r = p.add_run()
    _style_run(r, text, size, color, bold, font)
    return p

def rect(slide, l, t, w, h, fill=WHITE, line=BLACK, lw=1.5, rotation=0, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if round_:
        try: shp.adjustments[0] = 0.12
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    if rotation: shp.rotation = rotation
    return shp

def oline(slide, l, t, w, h, lw=2.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    shp.fill.background()
    shp.line.color.rgb = BLACK
    shp.line.width = Pt(lw)
    return shp

def sticker(slide, l, t, w, h, text, size=13, rotation=-3, fill=WHITE,
            color=BLACK, bold=False, font=F_KAI):
    shp = rect(slide, l, t, w, h, fill=fill, line=BLACK, lw=1.5, rotation=rotation)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _style_run(r, text, size, color, bold, font)
    return shp

def seal(slide, l, t, size, text, rotation=3):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, size, size)
    try: shp.adjustments[0] = 0.18
    except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = RED
    shp.line.color.rgb = BLACK; shp.line.width = Pt(1.25)
    shp.rotation = rotation
    tf = shp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('vert', 'eaVert')
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    size_in = size / 914400.0 if hasattr(size, 'emu') else float(size)
    font_pt = max(6, int(size_in * 72 / 3.4))
    _style_run(r, text, font_pt, WHITE, True, F_KAI)
    return shp

def cloud(slide, items, x0, y0, w, h, gapx=0.3, gapy=0.25, size=13, max_cols=4):
    x, y = x0, y0
    for i, t in enumerate(items):
        if i % max_cols == 0 and i > 0:
            x = x0; y += h + gapy
        sticker(slide, x, y, w, h, t, size=size, rotation=(3 if i % 2 else -2),
                fill=WHITE, color=BLACK, bold=True)
        x += w + gapx

def divider(idx, num_cn, en, title, quote, sticker_text):
    s = new_slide()
    # 水印
    tb(s, Inches(9.8), Inches(3.2), Inches(3.4), Inches(3.0), num_cn,
       size=200, color=YELLOW_D, bold=True, align=PP_ALIGN.RIGHT)
    page_chrome(s, idx)
    sticker(s, Inches(10.9), Inches(0.55), Inches(1.9), Inches(0.55),
            sticker_text, size=14, rotation=4, fill=RED, color=WHITE, bold=True)
    tb(s, Inches(0.9), Inches(1.5), Inches(4), Inches(2.0), num_cn,
       size=100, color=BLACK, bold=True)
    tb(s, Inches(0.9), Inches(3.25), Inches(10), Inches(1.2), title,
       size=58, color=BLACK, bold=True)
    rect(s, Inches(0.95), Inches(4.5), Inches(1.6), Inches(0.05), fill=BLACK, line=None)
    tb(s, Inches(0.9), Inches(4.75), Inches(10), Inches(0.7), quote,
       size=22, color=GRAY, font=F_KAI)
    return s

# ========== S1 封面 ==========
s = new_slide(); page_chrome(s, 0)
oline(s, Inches(9.2), Inches(-2.2), Inches(6.0), Inches(6.0), lw=2.5)
oline(s, Inches(10.4), Inches(-1.0), Inches(6.0), Inches(6.0), lw=1.25)
tb(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4),
   "献给所有不想被定义的年轻人", size=14, color=BLACK, bold=False)
sticker(s, Inches(10.35), Inches(0.5), Inches(2.5), Inches(0.5), "90后 · 00后 专场",
        size=13, rotation=3, fill=WHITE, bold=True)
tb_rich(s, Inches(1.0), Inches(2.15), Inches(11.5), Inches(2.2),
        [("人生是", 88, BLACK, True, F_SERIF), ("旷野", 92, RED, True, F_KAI)])
tb(s, Inches(1.05), Inches(4.15), Inches(10), Inches(0.7),
   "AI 时代 · 我们如何创造自己的世界", size=24, color=BLACK, font=F_KAI)
rect(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.03), fill=BLACK, line=None)
seal(s, Inches(0.9), Inches(6.2), Inches(0.85), "张医生", rotation=3)
tb(s, Inches(2.0), Inches(6.2), Inches(8), Inches(0.5),
   "AI 营养师 · 张医生", size=19, color=BLACK, bold=True)
tb(s, Inches(2.0), Inches(6.72), Inches(8), Inches(0.4),
   "医生 × 营养 × AI —— 一个很难被定义的人", size=12, color=GRAY, font=F_KAI)
tb(s, Inches(10.0), Inches(6.25), Inches(2.8), Inches(0.7),
   "不是轨道，是旷野。", size=20, color=BLACK, font=F_KAI,
   align=PP_ALIGN.RIGHT, rotation=-4)

# ========== S2 目录 ==========
s = new_slide(); page_chrome(s, 1)
tb(s, Inches(0.9), Inches(6.35), Inches(6), Inches(0.4),
   "翻开这一页，故事开始了 →", size=14, color=GRAY, font=F_KAI, rotation=-3)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "AGENDA", size=12, color=GRAY, font=F_SANS)
tb(s, Inches(0.9), Inches(1.0), Inches(10), Inches(1.0), "今天讲三件事", size=42, color=BLACK, bold=True)
cols = [
    ("01", "我的故事", "人生不设限 —— 一个医生、导演、音乐人、营养师、AI 玩家的多重身份。"),
    ("02", "世界的本质", "《人类简史》—— 规则、虚构与共同相信，世界由故事构成。"),
    ("03", "AI 时代", "知识门槛被砍平，读懂规则的人，开始创造自己的世界。"),
]
for i, (num, t, d) in enumerate(cols):
    x = Inches(0.9 + i * 4.1)
    tb(s, x, Inches(2.7), Inches(1.5), Inches(1.2), num, size=52, color=RED, bold=True)
    tb(s, x + Inches(0.02), Inches(3.7), Inches(3.6), Inches(0.6), t, size=24, color=BLACK, bold=True)
    rect(s, x + Inches(0.02), Inches(4.35), Inches(2.4), Inches(0.045), fill=BLACK, line=None)
    tb(s, x + Inches(0.02), Inches(4.6), Inches(3.7), Inches(1.6), d, size=13, color=GRAY, ls=1.4)

# ========== S3 篇章页 壹 ==========
divider(2, "壹", "ONE", "我的故事", "人生不设限 —— 先讲讲我是谁。", "个人经历")

# ========== S4 很难定义我 ==========
s = new_slide(); page_chrome(s, 3)
tb(s, Inches(0.9), Inches(6.35), Inches(5), Inches(0.4),
   "你，也是多重宇宙。", size=15, color=GRAY, font=F_KAI, rotation=4)
seal(s, Inches(11.55), Inches(5.55), Inches(0.8), "不设限", rotation=-3)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "MY IDENTITY", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.05), Inches(10), Inches(1.0),
        [("我，很难被", 44, BLACK, True, F_SERIF), ("定义", 46, RED, True, F_KAI)])
cloud(s, ["晚会导演", "主持人", "音乐", "旅行", "阅读", "医生",
          "营养师", "讲师", "AI"],
      Inches(0.9), Inches(2.35), Inches(2.7), Inches(0.65), gapx=0.35, gapy=0.3,
      size=15, max_cols=3)
tb_rich(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.2),
        [("每一个身份都是真实的，也都只是我的一部分。", 17, BLACK, True, F_SERIF),
         ("\n过去的我，做过太多「不像同一个人」的事 —— 站在舞台上，也进过手术室。", 14, GRAY, False, F_SERIF)])

# ========== S5 人生不设限 ==========
s = new_slide(); page_chrome(s, 4)
sticker(s, Inches(0.9), Inches(1.5), Inches(3.0), Inches(0.6), "我的第一个价值观",
        size=15, rotation=-2, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4),
        [("人生不设限。", 62, BLACK, True, F_SERIF)])
tb(s, Inches(0.92), Inches(3.75), Inches(11), Inches(0.6),
   "不要用传统制定的框架，去限定自己。", size=24, color=GRAY, font=F_KAI)
rect(s, Inches(0.95), Inches(5.1), Inches(2.0), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.9), Inches(5.4), Inches(11), Inches(0.5),
   "轨道是别人画好的；旷野，是自己走出来的。", size=16, color=BLACK)

# ========== S6 旷野 vs 轨道 ==========
s = new_slide(); page_chrome(s, 5)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "METAPHOR", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.0), Inches(11), Inches(1.0),
        [("人生是旷野，不是", 40, BLACK, True, F_SERIF), ("轨道", 42, RED, True, F_KAI)])
tb(s, Inches(9.8), Inches(6.3), Inches(3.0), Inches(0.5),
   "我选旷野。", size=18, color=GRAY, font=F_KAI, rotation=-4)
for i, (head, lines, rot) in enumerate([
        ("轨道", ["安全 · 确定 · 被安排", "父母画的线，社会铺的路", "平稳，但一眼望到头。"], 0),
        ("旷野", ["自由 · 未知 · 自己开路", "每一脚都可能踩空", "但每一步都通往真正的自己。"], 1.2)]):
    x = Inches(0.9 + i * 6.0)
    rect(s, x, Inches(2.35), Inches(5.5), Inches(3.3), fill=RGBColor(0xFF, 0xF7, 0xD1),
         line=BLACK, lw=2.5, rotation=rot)
    tb(s, x + Inches(0.4), Inches(2.75), Inches(4), Inches(0.7), head,
       size=28, color=BLACK, bold=True, rotation=rot)
    tb(s, x + Inches(0.4), Inches(3.7), Inches(4.8), Inches(1.6),
       "\n".join(lines), size=15, color=GRAY, ls=1.5, rotation=rot)

# ========== S7 写给零零后 ==========
s = new_slide(); page_chrome(s, 6)
sticker(s, Inches(0.9), Inches(1.4), Inches(2.4), Inches(0.6), "写给你们",
        size=15, rotation=2, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2),
        [("你是来", 44, BLACK, True, F_SERIF), ("创造规则", 48, RED, True, F_KAI),
         ("的，\n不是来遵守旧规则的。", 44, BLACK, True, F_SERIF)])
tb(s, Inches(0.92), Inches(4.7), Inches(11), Inches(0.6),
   "00 后，是定义新世界的一代。", size=22, color=GRAY, font=F_KAI)
tb(s, Inches(10.2), Inches(6.3), Inches(2.5), Inches(0.5),
   "去试，去错，去体验。", size=16, color=GRAY, font=F_KAI, rotation=5)

# ========== S8 篇章页 贰 ==========
divider(7, "贰", "TWO", "世界的本质", "《人类简史》 · 尤瓦尔·赫拉利 —— 世界的本质，是规则。", "一本书")

# ========== S9 一切都被规则规定 ==========
s = new_slide(); page_chrome(s, 8)
tb(s, Inches(0.9), Inches(6.35), Inches(6), Inches(0.4),
   "但规则，其实是「共同想象」。", size=14, color=GRAY, font=F_KAI, rotation=-3)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "RULES", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.6),
        [("你现在过的一切生活\n都是被", 36, BLACK, True, F_SERIF),
         ("规则", 38, RED, True, F_KAI),
         ("规定出来的", 36, BLACK, True, F_SERIF)])
cols = [
    ("父母的期望", "「为你好」三条路：稳定工作、早点结婚、别瞎折腾。"),
    ("学校的标准", "分数、排名、好专业 —— 一套统一的通关游戏。"),
    ("社会的剧本", "买房、还贷、升职、退休 —— 一条被写好的时间线。"),
]
for i, (t, d) in enumerate(cols):
    x = Inches(0.9 + i * 4.1)
    rect(s, x, Inches(3.35), Inches(3.7), Inches(2.6), fill=RGBColor(0xFF, 0xF7, 0xD1),
         line=BLACK, lw=1.75)
    tb(s, x + Inches(0.3), Inches(3.7), Inches(3.2), Inches(0.7), t,
       size=22, color=BLACK, bold=True)
    tb(s, x + Inches(0.3), Inches(4.6), Inches(3.1), Inches(1.2), d,
       size=13, color=GRAY, ls=1.4)

# ========== S10 钱是虚构的故事 ==========
s = new_slide(); page_chrome(s, 9)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "A TRUE STORY", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.0), Inches(6.5), Inches(1.6),
        [("钱，是虚构的", 40, BLACK, True, F_SERIF), ("故事", 42, RED, True, F_KAI)])
tb_rich(s, Inches(0.9), Inches(2.9), Inches(6.3), Inches(2.4),
        [("金钱、公司、国家、法律……", 17, BLACK, True, F_SERIF),
         ("\n这些都是人类共同相信的虚构。", 15, BLACK, False, F_SERIF),
         ("\n我们相信它，它就存在。", 15, BLACK, False, F_SERIF)])
sticker(s, Inches(0.9), Inches(5.35), Inches(2.6), Inches(0.6), "虚构 ≠ 谎言",
        size=15, rotation=2, fill=WHITE, bold=True)
# 右侧
sticker(s, Inches(8.5), Inches(2.3), Inches(3.8), Inches(1.5), "$  ¥  €",
        size=40, rotation=-2, fill=WHITE, color=RED, bold=True, font=F_SANS)
tb(s, Inches(8.2), Inches(4.35), Inches(4.4), Inches(2.2),
   "钞票本身只是纸，价值来自 80 亿人的共同相信 —— 这是人类协作的根基，也是每个人重新书写故事的权力。",
   size=14, color=GRAY, ls=1.5)

# ========== S11 人生的意义 ==========
s = new_slide(); page_chrome(s, 10)
sticker(s, Inches(0.9), Inches(1.4), Inches(3.6), Inches(0.6), "既然世界是虚构的",
        size=15, rotation=-2, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2),
        [("找到你真正想做的事，\n然后", 44, BLACK, True, F_SERIF),
         ("去探索、去体验", 48, RED, True, F_KAI),
         ("。", 44, BLACK, True, F_SERIF)])
tb(s, Inches(0.92), Inches(4.65), Inches(11), Inches(0.6),
   "体验，才是人生的意义。", size=22, color=GRAY, font=F_KAI)
tb(s, Inches(10.3), Inches(6.3), Inches(2.4), Inches(0.5),
   "世界由你重写。", size=16, color=GRAY, font=F_KAI, rotation=4)

# ========== S12 篇章页 叁 ==========
divider(11, "叁", "THREE", "AI 时代", "未来已来 —— 这不是预测，是现在。", "AI 时代")

# ========== S13 AI 创造一切 ==========
s = new_slide(); page_chrome(s, 12)
tb(s, Inches(0.9), Inches(6.35), Inches(6), Inches(0.4),
   "拐点已至。", size=14, color=GRAY, font=F_KAI, rotation=3)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "GENERATIVE", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.6),
        [("AI 可以创造\n一切你能", 36, BLACK, True, F_SERIF),
         ("想象", 38, RED, True, F_KAI),
         ("的东西", 36, BLACK, True, F_SERIF)])
cols = [
    ("90%", "的人无需重复劳动", "重复的、可复制的劳动，正在被 AI 接管。"),
    ("0", "知识门槛被砍平", "过去要靠十年积累的信息差，现在一句话就能调用。"),
    ("∞", "创造力的可能", "写作、设计、代码、视频 —— AI 放大每一个人的想法。"),
]
for i, (num, t, d) in enumerate(cols):
    x = Inches(0.9 + i * 4.1)
    tb(s, x, Inches(3.0), Inches(3.6), Inches(1.2), num, size=52, color=RED, bold=True)
    tb(s, x, Inches(4.1), Inches(3.6), Inches(0.6), t, size=20, color=BLACK, bold=True)
    rect(s, x, Inches(4.75), Inches(2.4), Inches(0.045), fill=BLACK, line=None)
    tb(s, x, Inches(4.95), Inches(3.7), Inches(1.4), d, size=13, color=GRAY, ls=1.4)

# ========== S14 规则向所有人打开 ==========
s = new_slide(); page_chrome(s, 13)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "THE NEW GAME", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.2),
        [("过去：信息差 = 特权。\n现在：AI 让信息差", 40, BLACK, True, F_SERIF),
         ("归零", 44, RED, True, F_KAI),
         ("。", 40, BLACK, True, F_SERIF)])
tb(s, Inches(0.92), Inches(4.6), Inches(11), Inches(0.7),
   "读懂规则的那一刻，就是你创造世界的那一刻。", size=24, color=GRAY, font=F_KAI)

# ========== S15 AI + 大健康 ==========
s = new_slide(); page_chrome(s, 14)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "MISSION", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.6),
        [("未来 AI 时代，\n做一件", 36, BLACK, True, F_SERIF),
         ("有价值", 38, RED, True, F_KAI),
         ("的事", 36, BLACK, True, F_SERIF)])
tb_rich(s, Inches(0.9), Inches(2.85), Inches(11), Inches(1.2),
        [("用 AI，把有价值的健康信息，传递给更多人。", 16, BLACK, True, F_SERIF),
         ("\n因为 ", 15, BLACK, False, F_SERIF),
         ("90% 的人还不知道", 17, RED, True, F_SERIF),
         ("：", 15, BLACK, False, F_SERIF)])
cloud(s, ["AI 的本质", "信息的本质", "健康的本质", "财富的本质"],
      Inches(0.9), Inches(4.35), Inches(2.8), Inches(0.7), gapx=0.4, gapy=0.3,
      size=16, max_cols=4)
tb(s, Inches(0.9), Inches(6.35), Inches(6), Inches(0.4),
   "知道规则的人，先到未来。", size=14, color=GRAY, font=F_KAI, rotation=-3)

# ========== S16 两大风口 ==========
s = new_slide(); page_chrome(s, 15)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "TWO WINDS", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.0), Inches(6), Inches(1.4),
        [("两大", 40, BLACK, True, F_SERIF), ("风口", 42, RED, True, F_KAI)])
tb(s, Inches(0.9), Inches(2.55), Inches(5.5), Inches(1.0),
   "风口不是等来的，是读懂规则后，站上去的。", size=15, color=GRAY, font=F_KAI, ls=1.4)
rect(s, Inches(0.9), Inches(3.7), Inches(5.3), Inches(2.4), fill=RGBColor(0xFF, 0xF7, 0xD1),
     line=BLACK, lw=2, rotation=0)
tb(s, Inches(1.25), Inches(4.0), Inches(4.6), Inches(0.6), "风口一 · AI", size=24, color=BLACK, bold=True)
tb(s, Inches(1.25), Inches(4.8), Inches(4.6), Inches(1.0),
   "先学会用 AI —— 让 AI 成为你的能力放大器。", size=14, color=GRAY, ls=1.4)
rect(s, Inches(6.7), Inches(3.7), Inches(5.6), Inches(2.4), fill=RGBColor(0xFF, 0xF7, 0xD1),
     line=BLACK, lw=2, rotation=1.2)
tb(s, Inches(7.05), Inches(4.0), Inches(4.9), Inches(0.6), "风口二 · 大健康", size=24, color=BLACK, bold=True)
tb(s, Inches(7.05), Inches(4.8), Inches(4.9), Inches(1.0),
   "把健康信息通过 AI 讲清楚 —— 帮更多人活得明白。", size=14, color=GRAY, ls=1.4)
sticker(s, Inches(0.9), Inches(6.45), Inches(4.4), Inches(0.6), "落地：短视频 × 自媒体",
        size=15, rotation=-2, fill=WHITE, bold=True)

# ========== S17 我在做的事 ==========
s = new_slide(); page_chrome(s, 16)
seal(s, Inches(11.55), Inches(5.5), Inches(0.8), "演示", rotation=-3)
tb(s, Inches(0.9), Inches(6.3), Inches(9), Inches(0.4),
   "（现场演示 → AI 工具 · 自媒体后台 · 知识库）", size=13, color=GRAY, font=F_KAI)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "WHAT I BUILD", size=12, color=GRAY, font=F_SANS)
tb(s, Inches(0.9), Inches(1.0), Inches(10), Inches(1.0), "我正在做的事", size=42, color=BLACK, bold=True)
cols = [
    ("AI", "AI 工具应用", "写作、做图、做视频、搭知识库 —— 一人即团队。"),
    ("DV", "大健康自媒体创作", "用短视频把健康的本质，讲给普通人听。"),
    ("KB", "本地知识库系统", "把十年的医生经验 + 营养认知，变成一套可复用的系统。"),
]
for i, (num, t, d) in enumerate(cols):
    x = Inches(0.9 + i * 4.1)
    rect(s, x, Inches(2.5), Inches(3.7), Inches(3.2), fill=RGBColor(0xFF, 0xF7, 0xD1),
         line=BLACK, lw=1.75)
    tb(s, x + Inches(0.3), Inches(2.85), Inches(3.1), Inches(0.9), num,
       size=34, color=RED, bold=True)
    tb(s, x + Inches(0.3), Inches(3.75), Inches(3.1), Inches(0.6), t,
       size=20, color=BLACK, bold=True)
    tb(s, x + Inches(0.3), Inches(4.5), Inches(3.1), Inches(1.0), d,
       size=13, color=GRAY, ls=1.4)

# ========== S18 结尾 ==========
s = new_slide(); page_chrome(s, 17)
sticker(s, Inches(0.9), Inches(1.3), Inches(3.0), Inches(0.6), "最后，你也可以",
        size=15, rotation=-2, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6),
        [("借 AI 这条通道，\n把你认为", 42, BLACK, True, F_SERIF),
         ("有价值", 46, RED, True, F_KAI),
         ("的东西，\n让更多人受益。", 42, BLACK, True, F_SERIF)])
tb_rich(s, Inches(0.92), Inches(5.05), Inches(11), Inches(0.8),
        [("AI 营养师 · 张医生 —— 谢谢，一起上路。", 20, GRAY, False, F_KAI)])
tb(s, Inches(10.0), Inches(6.25), Inches(2.8), Inches(0.6),
   "人生是旷野，不是轨道。", size=18, color=BLACK, font=F_KAI, align=PP_ALIGN.RIGHT, rotation=4)

# ========== SAVE ==========
out_dir = '/Users/mac/Documents/zfc最强大脑/A知识沉淀大脑/01-旧知识沉淀/01-营养知识'
out = os.path.join(out_dir, 'AI大健康-给年轻人的人生旷野课.pptx')
os.makedirs(out_dir, exist_ok=True)
prs.save(out)
print('PPT saved to:', out)
print('Total slides:', len(prs.slides))
