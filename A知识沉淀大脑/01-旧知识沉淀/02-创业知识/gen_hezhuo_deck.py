# -*- coding: utf-8 -*-
"""安利合作模式（OPP 大奖·三大奖金·龙卷风）· 黄底黑字杂志风 —— 由 KC-创业-009 知识卡渲染
运行：python3 gen_hezhuo_deck.py  → 直接在当前目录生成可打开编辑的 .pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

YELLOW = RGBColor(0xFF, 0xD6, 0x00)
YELLOW_D = RGBColor(0xF0, 0xC2, 0x00)
BLACK   = RGBColor(0x11, 0x11, 0x11)
RED     = RGBColor(0xD7, 0x00, 0x2F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x55, 0x55, 0x55)
F_SERIF = '宋体'
F_KAI   = '楷体'
F_SANS  = '黑体'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
N_TOTAL = 9   # 共 9 页

def _cjk(run, font_cn):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_cn)

def _style_run(r, text, size, color, bold, font):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    _cjk(r, font)

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = YELLOW
    return s

def page_chrome(s, idx):
    pbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(prs.slide_width * (idx+1) / N_TOTAL), Inches(0.045))
    pbar.fill.solid(); pbar.fill.fore_color.rgb = RED; pbar.line.fill.background()
    for (cx, cy, sx, sy) in [(0.12,0.12,1,1),(12.13,0.12,-1,1),(0.12,7.28,1,-1),(12.13,7.28,-1,-1)]:
        m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(0.16*sx), Inches(0.16))
        m.fill.solid(); m.fill.fore_color.rgb = BLACK; m.line.fill.background()
    tb(s, Inches(11.7), Inches(7.02), Inches(1.3), Inches(0.4),
       f"{idx+1:02d} — {N_TOTAL:02d}", size=11, color=BLACK, font=F_SANS, align=PP_ALIGN.RIGHT)

def tb(slide, l, t, w, h, text, size=18, color=BLACK, bold=False, font=F_SERIF,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, rotation=0, ls=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = anchor
    except Exception: pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    r = p.add_run()
    _style_run(r, text, size, color, bold, font)
    if rotation: box.rotation = rotation
    return box

def tb_rich(slide, l, t, w, h, parts, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, rotation=0, ls=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = anchor
    except Exception: pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    for (t, sz, c, b, f) in parts:
        r = p.add_run()
        _style_run(r, t, sz, c, b, f)
    if rotation: box.rotation = rotation
    return box

def rect(slide, l, t, w, h, fill=WHITE, line=BLACK, lw=1.5, rotation=0, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if round_:
        try: shp.adjustments[0] = 0.12
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    if rotation: shp.rotation = rotation
    return shp

def oline(slide, l, t, w, h, lw=2.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    shp.fill.background()
    shp.line.color.rgb = BLACK
    shp.line.width = Pt(lw)
    return shp

def sticker(slide, l, t, w, h, text, size=13, rotation=-3, fill=WHITE,
            color=BLACK, bold=False, font=F_KAI):
    shp = rect(slide, l, t, w, h, fill=fill, line=BLACK, lw=1.5, rotation=rotation)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _style_run(r, text, size, color, bold, font)
    return shp

# ========== S1 封面 ==========
s = new_slide(); page_chrome(s, 0)
oline(s, Inches(9.0), Inches(-2.4), Inches(6.2), Inches(6.2), lw=2.5)
oline(s, Inches(10.3), Inches(-1.1), Inches(6.2), Inches(6.2), lw=1.25)
tb(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4),
   "AI 营养师 · 张医生", size=14, color=BLACK)
sticker(s, Inches(10.5), Inches(0.5), Inches(2.2), Inches(0.55), "2026.09",
        size=13, rotation=3, fill=WHITE, bold=True)
tb_rich(s, Inches(1.0), Inches(2.1), Inches(12.2), Inches(2.3),
        [("安利", 96, BLACK, True, F_SERIF), ("合作模式", 92, RED, True, F_KAI)])
rect(s, Inches(1.05), Inches(4.45), Inches(2.4), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(1.05), Inches(4.7), Inches(11.5), Inches(0.7),
   "OPP 大奖 · 三大奖金分配制度 · 龙卷风商业模式", size=22, color=BLACK, font=F_KAI)
tb(s, Inches(1.0), Inches(6.15), Inches(11), Inches(0.6),
   "70 元 + 年满 22 岁，人人可为的大事业", size=20, color=GRAY, font=F_KAI, rotation=-3)

# ========== S2 为什么借安利这个平台 ==========
s = new_slide(); page_chrome(s, 1)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "WHY AMWAY", size=12, color=GRAY, font=F_SANS)
sticker(s, Inches(10.9), Inches(0.55), Inches(1.9), Inches(0.55), "为什么选它",
        size=14, rotation=4, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.8),
        [("为什么借", 44, BLACK, True, F_SERIF), ("安利", 48, RED, True, F_SERIF),
         ("这个平台？", 44, BLACK, True, F_SERIF)])
rect(s, Inches(0.95), Inches(3.7), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.95), Inches(4.05), Inches(11.5), Inches(3.0),
   "① 供应商支持：纽崔莱 93 年，稳定靠谱 · 商誉无风险（品牌不垮、产品不出事）\n"
   "② 人才体系支持：伙伴办卡挂你后面，开分店不用你发工资、不用你算税，公司直接各发各的\n"
   "③ 做大逻辑：一年 1000 万营业额，一个人干不动 —— 找 100 人每人干 10 万就轻松做大",
   size=20, color=BLACK, font=F_SANS, ls=1.4)

# ========== S3 合作条件 ==========
s = new_slide(); page_chrome(s, 2)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "准入条件", size=12, color=GRAY, font=F_SANS)
sticker(s, Inches(10.9), Inches(0.55), Inches(1.9), Inches(0.55), "人人可做",
        size=14, rotation=4, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.6),
        [("合作条件：", 40, BLACK, True, F_SERIF), ("只有两条", 44, RED, True, F_KAI)])
rect(s, Inches(0.95), Inches(3.5), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
# 两张条件卡
rect(s, Inches(0.95), Inches(3.9), Inches(5.4), Inches(2.3), fill=WHITE, line=BLACK, lw=2, round_=True)
tb(s, Inches(1.3), Inches(4.25), Inches(4.8), Inches(1.9),
   "① 年满 22 岁\n（学生不推荐办卡，守法律）", size=22, color=BLACK, font=F_SANS)
rect(s, Inches(7.0), Inches(3.9), Inches(5.4), Inches(2.3), fill=RED, line=BLACK, lw=2, round_=True)
tb(s, Inches(7.35), Inches(4.25), Inches(4.8), Inches(1.9),
   "② 70 元办卡\n（无囤货 · 无最低业绩 · 无加入费）", size=22, color=WHITE, font=F_SANS)
tb(s, Inches(0.95), Inches(6.45), Inches(11.5), Inches(0.9),
   "你是自己的经销商，直接跟公司签约 —— 不找你上线拿货、不受他业绩影响",
   size=18, color=GRAY, font=F_KAI)

# ========== S4 三大奖金 · 一、业绩奖金 ==========
s = new_slide(); page_chrome(s, 3)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "BONUS 01 · 业绩奖金", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.6),
        [("业绩奖金", 52, BLACK, True, F_SERIF), ("（自己做的钱）", 26, RED, True, F_KAI)])
rect(s, Inches(0.95), Inches(3.2), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.95), Inches(3.55), Inches(11.5), Inches(3.6),
   "· 一个陪跑方案约 2000 元/月，找 3 个客户 → 你的卡 6000 元营业额\n"
   "· 6000 × 3% = 180 元（另有 6% 电子券，可用于购物）\n"
   "· 做 30 个客户（6 万）→ 15% = 9000 元\n"
   "· 这个月干 0 就 0、不会被开除 —— 没有最低业绩、不要求囤货，自己是老板",
   size=20, color=BLACK, font=F_SANS, ls=1.4)

# ========== S5 三大奖金 · 二、市场开拓奖金 ==========
s = new_slide(); page_chrome(s, 4)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "BONUS 02 · 市场开拓奖金", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.5),
        [("市场开拓奖金", 48, BLACK, True, F_SERIF), ("（带团队 · 做大的关键）", 24, RED, True, F_KAI)])
rect(s, Inches(0.95), Inches(2.9), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.95), Inches(3.2), Inches(11.5), Inches(1.3),
   "你做 6000，再带 A/B/C/D 四人各自做 6000（这是他们自己的市场） → 整个部门 30000 × 12% = 3600",
   size=19, color=BLACK, font=F_SANS, ls=1.35)
rect(s, Inches(0.95), Inches(4.5), Inches(11.3), Inches(2.5), fill=WHITE, line=BLACK, lw=1.5, round_=True)
tb(s, Inches(1.3), Inches(4.75), Inches(10.7), Inches(2.2),
   "先发后面部门：4×180 + 自己的180 → 你拿 2880（180 业绩奖 + 2700 市场开拓奖）\n"
   "颠覆案例：市场最大的陈总拿 720，你带团队拿 2880 —— 市场最大的人收入不一定最高！\n"
   "公平三维度：多劳多得 · 可超越 · 助人自助（帮人，不是利用人）",
   size=18, color=BLACK, font=F_SANS, ls=1.35)

# ========== S6 三大奖金 · 三、6% 领导奖金 ==========
s = new_slide(); page_chrome(s, 5)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "BONUS 03 · 领导奖金", size=12, color=GRAY, font=F_SANS)
tb_rich(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.5),
        [("6% 领导奖金", 50, BLACK, True, F_SERIF), ("（资产保障 · 关键奖金）", 24, RED, True, F_KAI)])
rect(s, Inches(0.95), Inches(2.9), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.95), Inches(3.2), Inches(11.5), Inches(1.5),
   "有人把市场做到 12 万 5 → 你 12万5 × 6% = 7500 元/月 · 单独奖励给你，不影响他的收入\n"
   "保障 = 不做还有 —— 像一栋 5 层的房子月租 7500，而买这栋房要 300 万",
   size=19, color=BLACK, font=F_SANS, ls=1.35)
rect(s, Inches(0.95), Inches(4.8), Inches(11.3), Inches(2.3), fill=YELLOW_D, line=None, round_=True)
tb(s, Inches(1.3), Inches(5.0), Inches(10.7), Inches(2.0),
   "收入阶梯：1 组稳定 12万5（DD）≈ 年薪 30 万 → 3 个市场（营销经理）≈ 年薪百万\n"
   "→ 6 个市场（钻石）≈ 200 万+ → 20 个市场（FC）≈ 年入千万",
   size=20, color=BLACK, bold=True, font=F_SANS, ls=1.4)

# ========== S7 龙卷风商业模式 ==========
s = new_slide(); page_chrome(s, 6)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "龙卷风计划", size=12, color=GRAY, font=F_SANS)
sticker(s, Inches(10.9), Inches(0.55), Inches(1.9), Inches(0.55), "团队独有",
        size=14, rotation=4, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.6),
        [("龙卷风", 52, BLACK, True, F_SERIF), ("商业模式", 44, RED, True, F_KAI),
         ("（人脉众筹）", 24, BLACK, True, F_SANS)])
rect(s, Inches(0.95), Inches(3.2), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.95), Inches(3.55), Inches(11.5), Inches(3.6),
   "· 定义：人脉众筹 · 对等合作 · 共同发展 —— 把几路人脉凑到同一条线上\n"
   "· 人人都做「一得五」：市场共用，各自独立算钱，市场涨得飞快\n"
   "· 节奏：第一年卷第一条线到 12 万5 → 第二年卷第二条 → 第三年卷第三条\n"
   "· 三年百万年薪：三个 12万5 市场 = 营销经理（华工系统 · 林伟博士独创）",
   size=20, color=BLACK, font=F_SANS, ls=1.4)

# ========== S8 市场空间 ==========
s = new_slide(); page_chrome(s, 7)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "MARKET SIZE", size=12, color=GRAY, font=F_SANS)
sticker(s, Inches(10.9), Inches(0.55), Inches(1.9), Inches(0.55), "大事业",
        size=14, rotation=4, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.6),
        [("为什么是", 40, BLACK, True, F_SERIF), ("大事业", 50, RED, True, F_KAI)])
rect(s, Inches(0.95), Inches(3.3), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.95), Inches(3.65), Inches(11.5), Inches(3.4),
   "· 1000 万人办过卡，真正做的只有 1%（约 10 万人）—— 中国 14 亿人口，市场远没开始\n"
   "· 全球 100 多个国家和地区，奖金制度、纽崔莱全球统一、健康相通\n"
   "· 有海外朋友就能做国际市场",
   size=20, color=BLACK, font=F_SANS, ls=1.4)

# ========== S9 结尾 / 金句 ==========
s = new_slide(); page_chrome(s, 8)
tb(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.4), "SUMMARY", size=12, color=GRAY, font=F_SANS)
seal = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.5), Inches(5.0), Inches(0.9), Inches(0.9))
seal.fill.solid(); seal.fill.fore_color.rgb = RED; seal.line.color.rgb = BLACK; seal.line.width = Pt(1.25)
seal.rotation = 3
tf = seal.text_frame; tf.word_wrap = False; tf.auto_size = MSO_AUTO_SIZE.NONE
tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
bodyPr = tf._txBody.find(qn('a:bodyPr')); bodyPr.set('vert', 'eaVert')
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); _style_run(r, "共赢", 15, WHITE, True, F_KAI)
tb_rich(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.2),
        [("安利的公平，", 40, BLACK, True, F_SERIF),
         ("是算账算出来的", 44, RED, True, F_KAI)])
rect(s, Inches(0.95), Inches(3.9), Inches(2.2), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.95), Inches(4.25), Inches(10.5), Inches(2.2),
   "· 多劳多得 · 可超越 · 助人自助 · 非零和\n"
   "· 先发后面部门的钱，谁都不被辜负\n"
   "· 70 元 + 22 岁，普通人也能做的大事业，一条线即可退休",
   size=20, color=BLACK, font=F_SANS, ls=1.4)

out_dir = '/Users/mac/Documents/zfc最强大脑/A知识沉淀大脑/01-旧知识沉淀/02-创业知识'
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, '安利合作模式-OPP大奖·三大奖金·龙卷风.pptx')
prs.save(out)
print('PPT saved to:', out)
print('Total slides:', len(prs.slides))
