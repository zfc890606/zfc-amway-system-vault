# -*- coding: utf-8 -*-
"""安利合作模式 · OPP 三大奖金分配制度 —— 专业蓝金风（16:9）
内容源：A知识沉淀大脑/01-旧知识沉淀/02-创业知识/KC-创业-009-安利合作模式OPP三大奖金分配制度.md
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 主题色 ----------
NAVY   = RGBColor(0x0F, 0x2A, 0x4E)   # 深海军蓝
NAVY2  = RGBColor(0x1B, 0x3A, 0x63)   # 次深蓝
GOLD   = RGBColor(0xC9, 0xA2, 0x27)   # 金色
GOLD_L = RGBColor(0xF0, 0xDD, 0xA8)   # 浅金
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x22, 0x22, 0x22)
GRAY   = RGBColor(0x66, 0x66, 0x66)
RED    = RGBColor(0xB4, 0x3A, 0x2A)   # 强调红
GREEN  = RGBColor(0x2E, 0x7D, 0x5B)   # 强调绿
LIGHT  = RGBColor(0xF3, 0xF5, 0xF9)   # 浅灰底
LINEC  = RGBColor(0xD5, 0xDC, 0xE6)   # 浅边框

F_HEAD = '黑体'
F_BODY = '宋体'
F_KAI  = '楷体'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------- 工具 ----------
def _cjk(run, font_cn):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_cn)

def style_run(r, text, size, color, bold, font):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    _cjk(r, font)

def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def rect(slide, l, t, w, h, fill, line=None, lw=1.0, round_=False, adj=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if round_:
        try: shp.adjustments[0] = adj
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def textbox(slide, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            ls=1.15, space_after=6):
    """paras: list of paragraphs; each = list of (text,size,color,bold,font) runs."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = anchor
    except Exception: pass
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        p.space_after = Pt(space_after)
        for (txt, sz, c, b, f) in runs:
            r = p.add_run()
            style_run(r, txt, sz, c, b, f)
    return box

def T(slide, l, t, w, h, text, size=16, color=BLACK, bold=False, font=F_BODY,
      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.15):
    return textbox(slide, l, t, w, h, [[(text, size, color, bold, font)]],
                   align=align, anchor=anchor, ls=ls)

def title_bar(slide, title, kicker=None):
    """内容页标准标题：海军蓝条 + 金色下划线 + 右上角页码锚点"""
    rect(slide, 0, 0, prs.slide_width, Inches(0.16), NAVY)
    rect(slide, Inches(0.55), Inches(0.62), Inches(0.09), Inches(0.62), GOLD)
    T(slide, Inches(0.82), Inches(0.55), Inches(11.0), Inches(0.8), title,
      size=28, color=NAVY, bold=True, font=F_HEAD)
    if kicker:
        T(slide, Inches(0.84), Inches(1.18), Inches(11.0), Inches(0.4), kicker,
          size=12, color=GRAY, font=F_BODY)

def chip(slide, l, t, w, h, text, fill=NAVY2, color=WHITE, size=14, bold=True, font=F_HEAD):
    shp = rect(slide, l, t, w, h, fill, round_=True, adj=0.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); style_run(r, text, size, color, bold, font)
    return shp

def set_cell(cell, runs_list, fill=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = anchor
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for (txt, sz, c, b, f) in runs_list:
        r = p.add_run()
        style_run(r, txt, sz, c, b, f)

def add_table(slide, l, t, w, h, rows_data, col_widths, header_fill=NAVY,
              header_color=WHITE, row_fill=None, body_size=14, header_size=15):
    """rows_data: list of lists; first row = header."""
    n_rows = len(rows_data); n_cols = len(rows_data[0])
    shape = slide.shapes.add_table(n_rows, n_cols, l, t, w, h)
    table = shape.table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw
    table.rows[0].height = Inches(0.5)
    for ri in range(1, n_rows):
        table.rows[ri].height = Inches(0.45)
    for ri, row in enumerate(rows_data):
        for ci, cellval in enumerate(row):
            cell = table.cell(ri, ci)
            if ri == 0:
                set_cell(cell, [(str(cellval), header_size, header_color, True, F_HEAD)],
                         fill=header_fill, align=PP_ALIGN.CENTER)
            else:
                fill = LIGHT if (row_fill is None and ri % 2 == 0) else row_fill
                set_cell(cell, [(str(cellval), body_size, BLACK, False, F_BODY)], fill=fill)
    return table

def footer(slide, idx, total=14):
    rect(slide, 0, Inches(7.28), prs.slide_width, Inches(0.22), NAVY)
    T(slide, Inches(0.4), Inches(7.28), Inches(8), Inches(0.2),
      "AI营养师 · 张医生  |  安利合作模式 OPP 三大奖金", size=8, color=GOLD_L, font=F_BODY)
    T(slide, Inches(11.6), Inches(7.28), Inches(1.4), Inches(0.2),
      f"{idx:02d} / {total:02d}", size=8, color=GOLD_L, font=F_BODY, align=PP_ALIGN.RIGHT)

# ==================================================================
# S1 封面
# ==================================================================
s = new_slide(NAVY)
rect(s, 0, 0, prs.slide_width, Inches(0.18), GOLD)
rect(s, Inches(1.0), Inches(2.02), Inches(0.9), Inches(0.06), GOLD)
T(s, Inches(1.0), Inches(0.85), Inches(11.0), Inches(0.6),
  "AI营养师 · 张医生", size=16, color=GOLD_L, font=F_HEAD)
T(s, Inches(1.0), Inches(2.35), Inches(11.5), Inches(1.5),
  "安利合作模式", size=64, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(1.0), Inches(3.55), Inches(11.5), Inches(0.7),
  "OPP 事业机会 × 三大奖金分配制度", size=24, color=GOLD_L, font=F_HEAD)
T(s, Inches(1.0), Inches(4.6), Inches(11.5), Inches(0.9),
  "70元 · 22岁 · 人人可为\n多劳多得 · 可超越 · 助人自助 · 非零和", size=16, color=WHITE, font=F_BODY, ls=1.4)
rect(s, Inches(1.0), Inches(6.15), Inches(3.0), Inches(0.03), GOLD)
T(s, Inches(1.0), Inches(6.4), Inches(11.5), Inches(0.5),
  "2026-09-01  ·  根据张医生录音《大健康与安利合作模式及AI应用分享》整理", size=12, color=GOLD_L, font=F_BODY)

# ==================================================================
# S2 目录
# ==================================================================
s = new_slide(); title_bar(s, "目录", "安利合作模式 · 一次讲透")
items = [
    ("01", "为什么借安利平台", "供应商 + 人才体系两大支持"),
    ("02", "合作条件 · 人人可为", "70元 · 22岁 · 无囤货无业绩"),
    ("03", "三大奖金机制", "业绩奖金 · 市场开拓 · 6%领导奖"),
    ("04", "7人算账 · 公平是算出来的", "市场最大 ≠ 收入最高"),
    ("05", "收入阶梯 · 市场=资产", "DD → 营销经理 → 钻石 → FC"),
    ("06", "龙卷风 · 三年百万年薪", "人脉众筹 · 一年一条线"),
]
for i, (num, main, sub) in enumerate(items):
    col = i // 3; row = i % 3
    l = Inches(0.9 + col * 6.05); t = Inches(1.8 + row * 1.72)
    rect(s, l, t, Inches(5.6), Inches(1.45), LIGHT, round_=True, adj=0.12)
    rect(s, l, t, Inches(0.09), Inches(1.45), GOLD)
    T(s, Inches(l + 0.28), t + Inches(0.14), Inches(1.2), Inches(1.0),
      num, size=30, color=GOLD, bold=True, font=F_HEAD)
    T(s, Inches(l + 0.95), t + Inches(0.16), Inches(4.4), Inches(0.55),
      main, size=19, color=NAVY, bold=True, font=F_HEAD)
    T(s, Inches(l + 0.95), t + Inches(0.82), Inches(4.5), Inches(0.5),
      sub, size=12, color=GRAY, font=F_BODY)
footer(s, 2)

# ==================================================================
# S3 大框架
# ==================================================================
s = new_slide(); title_bar(s, "先看大框架：互联网 × 大健康", "不是卖产品，卖产品没有前途——是卖解决方案")
# 流程条
rect(s, Inches(0.9), Inches(1.75), Inches(4.1), Inches(1.35), NAVY, round_=True)
T(s, Inches(1.15), Inches(1.95), Inches(3.7), Inches(0.5), "互联网 · 公域获客", size=18, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(1.15), Inches(2.5), Inches(3.7), Inches(0.5), "社群运营 · 公域流量 · 引流进来", size=12, color=GOLD_L, font=F_BODY)
# 箭头
T(s, Inches(5.15), Inches(2.0), Inches(1.0), Inches(0.8), "→", size=40, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
rect(s, Inches(6.3), Inches(1.75), Inches(6.1), Inches(1.35), NAVY2, round_=True)
T(s, Inches(6.55), Inches(1.95), Inches(5.6), Inches(0.5), "大健康 · 解决方案落地", size=18, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(6.55), Inches(2.5), Inches(5.7), Inches(0.5), "慢病逆转 · 肥胖/痛风/糖尿病/高血压 · 健康人群落地", size=12, color=GOLD_L, font=F_BODY)
# 两大支持
T(s, Inches(0.9), Inches(3.55), Inches(6.0), Inches(0.5), "安利在这个框架里，帮我们什么？", size=16, color=NAVY, bold=True, font=F_HEAD)
for i, (t1, t2, t3) in enumerate([
    ("① 供应商支持", "纽崔莱93年 · 稳定靠谱 · 商誉无风险", "品牌倒了白干、产品出事人品都没了——选放心的"),
    ("② 人才平台支持", "伙伴办卡挂你后面 · 开分店不用你发工资", "不用开公司、不用算税，公司给各人发钱——你只管获客+带人"),
]):
    l = Inches(0.9 + i * 6.05)
    rect(s, l, Inches(4.2), Inches(5.6), Inches(1.5), LIGHT, round_=True)
    T(s, l + Inches(0.25), Inches(4.4), Inches(5.1), Inches(0.5), t1, size=17, color=NAVY, bold=True, font=F_HEAD)
    T(s, l + Inches(0.25), Inches(4.95), Inches(5.1), Inches(0.6), t2, size=13, color=BLACK, font=F_BODY, ls=1.2)
    T(s, l + Inches(0.25), Inches(5.5), Inches(5.1), Inches(0.5), t3, size=11, color=GRAY, font=F_BODY, ls=1.1)
# 做大逻辑
rect(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.95), NAVY, round_=True)
T(s, Inches(1.15), Inches(6.16), Inches(11.0), Inches(0.6),
  "做大的逻辑：一年1000万营业额，一个人干不动 → 找100个人，每人干10万，就轻松做大",
  size=15, color=WHITE, bold=True, font=F_BODY)
footer(s, 3)

# ==================================================================
# S4 合作条件
# ==================================================================
s = new_slide(); title_bar(s, "合作条件：人人可为", "两个条件而已，不要求囤货、没有业绩要求、没有加入费")
for i, (num, main, sub) in enumerate([
    ("①", "年满 22 岁", "学生不推荐办卡（遵守国家法律，学习为主）"),
    ("②", "70 元办卡", "门槛低到每个人都能起步"),
]):
    l = Inches(0.9 + i * 6.05)
    rect(s, l, Inches(1.75), Inches(5.6), Inches(1.9), NAVY, round_=True)
    T(s, l + Inches(0.3), Inches(2.0), Inches(1.0), Inches(0.8), num, size=40, color=GOLD, bold=True, font=F_HEAD)
    T(s, l + Inches(1.15), Inches(2.15), Inches(4.3), Inches(0.7), main, size=26, color=WHITE, bold=True, font=F_HEAD)
    T(s, l + Inches(1.15), Inches(2.85), Inches(4.3), Inches(0.6), sub, size=12, color=GOLD_L, font=F_BODY)
# 三个无
for i, t in enumerate(["不要求囤货", "无最低业绩要求", "无加入费"]):
    chip(s, Inches(0.9 + i * 2.15), Inches(3.95), Inches(1.95), Inches(0.55), t, fill=NAVY2, size=15)
# 对比框
rect(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.9), LIGHT, round_=True)
T(s, Inches(1.15), Inches(5.05), Inches(11.0), Inches(0.5), "对比：传统生意合作", size=16, color=NAVY, bold=True, font=F_HEAD)
T(s, Inches(1.15), Inches(5.55), Inches(11.0), Inches(0.9),
  "红牛代理要 300 万才跟你合作；很多平台要交 2-3 万加入费。\n安利：70 块就能开张——你是自己的老板，跟公司签约，不是给上线打工、不从上线拿货。",
  size=14, color=BLACK, font=F_BODY, ls=1.35)
footer(s, 4)

# ==================================================================
# S5 三大奖金总览
# ==================================================================
s = new_slide(); title_bar(s, "三大奖金机制", "自己做 → 带团队 → 不做也有，三级递进")
cards = [
    ("①", "业绩奖金", "自己做", "你干多少、拿多少", "3% ~ 21% 月结", NAVY),
    ("②", "市场开拓奖金", "带团队", "开分店，做大靠复制", "部门差额归你", NAVY2),
    ("③", "6% 领导奖", "不做也有", "资产保障·可躺平", "12.5万×6%/月", GOLD),
]
for i, (num, name, tag, sub, note, fill) in enumerate(cards):
    l = Inches(0.9 + i * 4.05)
    rect(s, l, Inches(2.0), Inches(3.7), Inches(4.2), LIGHT, round_=True)
    rect(s, l, Inches(2.0), Inches(3.7), Inches(1.1), fill, round_=True, adj=0.18)
    T(s, l + Inches(0.25), Inches(2.12), Inches(1.0), Inches(0.8), num, size=28, color=WHITE, bold=True, font=F_HEAD)
    T(s, l + Inches(0.9), Inches(2.18), Inches(2.6), Inches(0.7), name, size=22, color=WHITE, bold=True, font=F_HEAD)
    chip(s, l + Inches(0.25), Inches(3.35), Inches(3.2), Inches(0.5), tag, fill=NAVY2 if fill != GOLD else GOLD, size=14)
    T(s, l + Inches(0.3), Inches(4.1), Inches(3.2), Inches(0.6), sub, size=15, color=BLACK, bold=True, font=F_HEAD)
    T(s, l + Inches(0.3), Inches(4.75), Inches(3.2), Inches(0.6), note, size=13, color=GRAY, font=F_BODY)
T(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
  "讲制度不讲复杂比例表（9种12项·哈佛都在学），只讲这三样，普通人就能看懂。",
  size=13, color=GRAY, font=F_KAI)
footer(s, 5)

# ==================================================================
# S6 ① 业绩奖金
# ==================================================================
s = new_slide(); title_bar(s, "① 业绩奖金：自己做的钱", "月月结算，对照同一张比例表，自己做多做少自己决定")
# 两个算账卡
for i, (n, calc, res) in enumerate([
    ("例1", "3个客户 × 2000元/月 = 6000", "6000 × 3% = 180 元  (+6%电子券360)"),
    ("例2", "30个客户 = 6万元营业额", "6万 × 15% = 9000 元"),
]):
    l = Inches(0.9 + i * 6.05)
    rect(s, l, Inches(1.8), Inches(5.6), Inches(1.7), LIGHT, round_=True)
    chip(s, l + Inches(0.25), Inches(2.0), Inches(1.0), Inches(0.5), n, fill=GOLD, size=15)
    T(s, l + Inches(0.25), Inches(2.65), Inches(5.1), Inches(0.5), calc, size=15, color=BLACK, bold=True, font=F_BODY)
    T(s, l + Inches(0.25), Inches(3.15), Inches(5.1), Inches(0.5), res, size=16, color=RED, bold=True, font=F_BODY)
# 要点
rect(s, Inches(0.9), Inches(3.85), Inches(11.5), Inches(1.35), NAVY, round_=True)
T(s, Inches(1.15), Inches(4.0), Inches(11.0), Inches(1.0),
  "这个月干0，就0收入——但不会被开除。\n没有最低业绩要求、不要求囤货，自由企业：你是自己的老板，它是你的供应商。",
  size=15, color=WHITE, font=F_BODY, ls=1.3)
# 表
T(s, Inches(0.9), Inches(5.5), Inches(8), Inches(0.4), "比例表速查（BV 净营业额档位）", size=14, color=NAVY, bold=True, font=F_HEAD)
rows = [
    ["净营业额", "7500", "25000", "50000", "87500", "125000"],
    ["奖金比率", "6%", "12%", "15%", "18%", "21%"],
]
table = add_table(s, Inches(0.9), Inches(5.95), Inches(8.5), Inches(0.95), rows,
                  [Inches(1.7)] + [Inches(1.36)] * 5, body_size=13, header_size=13)
footer(s, 6)

# ==================================================================
# S7 ② 市场开拓奖金
# ==================================================================
s = new_slide(); title_bar(s, "② 市场开拓奖金：带团队的钱", "你不必很厉害，只需复制能力——把6000的能力教会给别人")
# 步骤条
rect(s, Inches(0.9), Inches(1.75), Inches(3.6), Inches(1.3), NAVY, round_=True)
T(s, Inches(1.1), Inches(1.9), Inches(3.2), Inches(0.5), "你：做 6000", size=17, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(1.1), Inches(2.45), Inches(3.2), Inches(0.5), "3个客户 × 2000/月", size=12, color=GOLD_L, font=F_BODY)
rect(s, Inches(4.75), Inches(1.75), Inches(3.7), Inches(1.3), NAVY2, round_=True)
T(s, Inches(4.95), Inches(1.9), Inches(3.3), Inches(0.5), "再带 A/B/C/D", size=17, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(4.95), Inches(2.45), Inches(3.3), Inches(0.5), "教会他们各自也做6000", size=12, color=GOLD_L, font=F_BODY)
rect(s, Inches(8.7), Inches(1.75), Inches(3.7), Inches(1.3), GOLD, round_=True)
T(s, Inches(8.9), Inches(1.9), Inches(3.3), Inches(0.5), "部门 = 5×6000 = 3万", size=17, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(8.9), Inches(2.45), Inches(3.3), Inches(0.5), "部门奖金 12% = 3600", size=12, color=WHITE, font=F_BODY)
# 算账
rect(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(2.0), LIGHT, round_=True)
T(s, Inches(1.15), Inches(3.55), Inches(11.0), Inches(0.5), "公司怎么分这笔钱？", size=15, color=NAVY, bold=True, font=F_HEAD)
T(s, Inches(1.15), Inches(4.05), Inches(11.0), Inches(1.1),
  "① 先发后面的部门：4个新人各拿走 180（4×180=720）——不能少别人的钱。\n② 再算你：部门奖金3600 − 720 − 你也是6000的业绩奖180 → 你拿 2880。\n③ 拆解：2880 = 180（自己的业绩奖） + 2700（市场开拓奖）",
  size=15, color=BLACK, font=F_BODY, ls=1.4)
T(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.8),
  "关键：那6000是他们自己做的、他们的市场，你只是教会了他们——你的收入来自“付出4份劳动”，不是剥削他们。",
  size=13, color=RED, bold=True, font=F_KAI)
footer(s, 7)

# ==================================================================
# S8 7人算账
# ==================================================================
s = new_slide(); title_bar(s, "7人算账：公平是算出来的", "讲给伙伴听的黄金故事——同一个部门，7个人，钱怎么分")
rows = [
    ["角色", "部门业绩", "这个月收入", "说明"],
    ["陈总（市场最大）", "7×6000 = 42000", "720 元", "市场最大，但收入不是最高"],
    ["张医生（带1人）", "6×6000 = 36000", "720 元", "只付出1份劳动"],
    ["你（带4人）", "5×6000 = 30000", "2880 元", "付出4份劳动，多劳多得"],
    ["4位新人", "各 6000", "各 180 元", "干多少拿多少"],
]
add_table(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(2.6), rows,
          [Inches(3.0), Inches(3.0), Inches(2.5), Inches(3.0)],
          body_size=14, header_size=14)
# 结论
rect(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.0), NAVY, round_=True)
T(s, Inches(1.15), Inches(5.05), Inches(11.0), Inches(0.6),
  "市场最大的人，收入不一定最高 —— 系统只看劳动量，不看资历。",
  size=17, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.8),
  "更颠覆：陈总这个月躺平干0 → 先拨完后面部门的钱后，他自己 0 元。钱先发后面的部门，谁都不被辜负。",
  size=14, color=BLACK, font=F_BODY)
footer(s, 8)

# ==================================================================
# S9 公平三维度
# ==================================================================
s = new_slide(); title_bar(s, "制度为什么留得住人", "四大特性——跟传统行业对比着讲")
cards = [
    ("多劳多得", "你拿2880，因为你付出4份劳动（带4个人）；我只带1个，只付1份。系统按劳动量算，不是按关系算。", NAVY),
    ("可超越", "我加入比陈总晚，能力上来就能超越——大家都参照同一张表，没有例外。传统要看领导调不调政策。", NAVY2),
    ("助人自助", "你不帮A/B/C/D，他们都做0，你做6000还是180；帮他们赚到钱，你的收入才上来。是帮人，不是利用人。", NAVY2),
    ("非零和", "传统是分蛋糕（你多他少、主任只有一个、争破头）；这里做多少拿多少，主任没有人数限制，一万个主任都有。", NAVY),
]
for i, (t1, t2, fill) in enumerate(cards):
    col = i // 2; row = i % 2
    l = Inches(0.9 + col * 6.05); t = Inches(1.85 + row * 2.5)
    rect(s, l, t, Inches(5.6), Inches(2.25), LIGHT, round_=True)
    rect(s, l, t, Inches(0.09), Inches(2.25), GOLD)
    T(s, l + Inches(0.3), t + Inches(0.18), Inches(5.0), Inches(0.6), t1, size=20, color=NAVY, bold=True, font=F_HEAD)
    T(s, l + Inches(0.3), t + Inches(0.85), Inches(5.05), Inches(1.3), t2, size=13, color=BLACK, font=F_BODY, ls=1.25)
T(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.5),
  "我影响不了你的收入，你也影响不了我的——每个人都是跟公司签约的独立经销商。",
  size=13, color=GRAY, font=F_KAI)
footer(s, 9)

# ==================================================================
# S10 ③ 6%领导奖
# ==================================================================
s = new_slide(); title_bar(s, "③ 6% 领导奖：不做也有", "这是张医生辞职做安利的核心原因——资产性收入")
# 逻辑流
rect(s, Inches(0.9), Inches(1.75), Inches(3.7), Inches(1.3), NAVY, round_=True)
T(s, Inches(1.1), Inches(1.9), Inches(3.3), Inches(0.6), "你培养的人做到 12.5万", size=15, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(1.1), Inches(2.5), Inches(3.3), Inches(0.5), "市场到21%封顶，你没差额了", size=11, color=GOLD_L, font=F_BODY)
rect(s, Inches(4.85), Inches(1.75), Inches(3.7), Inches(1.3), NAVY2, round_=True)
T(s, Inches(5.05), Inches(1.9), Inches(3.3), Inches(0.6), "公司额外奖励 6%", size=15, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(5.05), Inches(2.5), Inches(3.3), Inches(0.5), "不影响他的收入", size=11, color=GOLD_L, font=F_BODY)
rect(s, Inches(8.8), Inches(1.75), Inches(3.6), Inches(1.3), GOLD, round_=True)
T(s, Inches(9.0), Inches(1.9), Inches(3.2), Inches(0.6), "12.5万 × 6% = 7500/月", size=17, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(9.0), Inches(2.5), Inches(3.2), Inches(0.5), "市场稳健，不做也有", size=12, color=WHITE, font=F_BODY)
# 房子比喻
rect(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(2.1), LIGHT, round_=True)
T(s, Inches(1.15), Inches(3.55), Inches(11.0), Inches(0.5), "房子比喻：收入 = 收租", size=16, color=NAVY, bold=True, font=F_HEAD)
T(s, Inches(1.15), Inches(4.1), Inches(11.0), Inches(1.1),
  "在佛山有一栋5层楼，每层月租1500 → 一个月收租 7500。\n买下这栋楼要 300 万；而在安利，通过大健康做一组稳健市场，一年不到就能拿到这 7500/月，还不用背房贷。",
  size=15, color=BLACK, font=F_BODY, ls=1.4)
T(s, Inches(1.15), Inches(5.25), Inches(11.0), Inches(0.5),
  "保障 = 拥有资产：停下来也有钱，像房子一样。2020年到现在，睡到自然醒、不打卡。", size=13, color=RED, bold=True, font=F_KAI)
footer(s, 10)

# ==================================================================
# S11 收入阶梯
# ==================================================================
s = new_slide(); title_bar(s, "收入阶梯：市场 = 资产", "数字是示例教学，实际以公司当期制度为准")
rows = [
    ["达标", "奖衔", "年收入", "特点"],
    ["1组稳定 12.5万", "DD 高级主任", "30-40 万", "不做还有 · 可退休"],
    ["3个DD市场", "营销经理", "100 万+", "三条线不做还有"],
    ["6个DD市场", "钻石", "200 万+", "收入再上一个台阶"],
    ["20个DD市场", "FC 皇冠大使", "1000 万+", "张医生的目标"],
]
add_table(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.4), rows,
          [Inches(3.4), Inches(3.0), Inches(2.6), Inches(2.5)],
          body_size=15, header_size=15)
T(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.8),
  "35岁以后这些收入“不做还有”——这正是资产性收入与打工的区别：传统要干到60岁才停得下来。",
  size=15, color=NAVY, bold=True, font=F_HEAD)
T(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
  "张医生现状：1组市场，30-40万/年；正在卷第2、3组 → 目标3条线年薪百万。", size=13, color=GRAY, font=F_BODY)
footer(s, 11)

# ==================================================================
# S12 龙卷风
# ==================================================================
s = new_slide(); title_bar(s, "龙卷风商业模式：三年百万年薪", "人脉众筹 · 对等合作 · 共同发展（林伟博士/华工系统独有）")
# 定义
rect(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(1.3), NAVY, round_=True)
T(s, Inches(1.15), Inches(1.9), Inches(11.0), Inches(0.9),
  "一个人做市场，靠的是自己的资源和人脉；龙卷风 = 把人脉众筹到同一条线上。\n市场共用、各自独立算钱——人人都做“一得五”，市场涨得飞快。",
  size=16, color=WHITE, font=F_BODY, ls=1.35)
# 节奏三卡
for i, (n, t1, t2) in enumerate([
    ("第一年", "卷第一条线", "把一个市场做到 21% / 12.5万"),
    ("第二年", "卷第二条线", "再找三五人，龙卷风再来一条"),
    ("第三年", "卷第三条线", "三条线 = 营销经理 年薪百万"),
]):
    l = Inches(0.9 + i * 4.05)
    rect(s, l, Inches(3.4), Inches(3.7), Inches(1.9), LIGHT, round_=True)
    chip(s, l + Inches(0.25), Inches(3.6), Inches(3.2), Inches(0.5), n, fill=GOLD if i == 2 else NAVY2, size=15)
    T(s, l + Inches(0.28), Inches(4.25), Inches(3.2), Inches(0.55), t1, size=18, color=NAVY, bold=True, font=F_HEAD)
    T(s, l + Inches(0.28), Inches(4.8), Inches(3.2), Inches(0.5), t2, size=12, color=GRAY, font=F_BODY, ls=1.15)
T(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.8),
  "为什么是三条线、每人贡献三个人脉？——林伟博士（统计学背景）设计，全安利只有华工系统独有。",
  size=14, color=BLACK, font=F_BODY)
T(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6),
  "附加福利：达标还有海外旅游（张医生已带爸爸免费去过多个国家）。", size=13, color=GRAY, font=F_KAI)
footer(s, 12)

# ==================================================================
# S13 市场空间
# ==================================================================
s = new_slide(); title_bar(s, "为什么是大事业", "市场远没开始 + 全球统一制度 + 洗牌红利")
big = [
    ("1000万人", "办过卡", "真正在做的只有 1%"),
    ("14亿", "中国人口", "市场远没开始"),
    ("100+", "国家/地区", "制度统一·纽崔莱统一·健康相通"),
]
for i, (n, u, sub) in enumerate(big):
    l = Inches(0.9 + i * 4.05)
    rect(s, l, Inches(1.9), Inches(3.7), Inches(2.4), LIGHT, round_=True)
    T(s, l + Inches(0.3), Inches(2.1), Inches(3.2), Inches(1.0), n, size=34, color=NAVY, bold=True, font=F_HEAD)
    T(s, l + Inches(0.3), Inches(3.15), Inches(3.2), Inches(0.5), u, size=14, color=GOLD, bold=True, font=F_HEAD)
    T(s, l + Inches(0.3), Inches(3.7), Inches(3.2), Inches(0.5), sub, size=12, color=GRAY, font=F_BODY)
# 洗牌红利
rect(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.9), NAVY, round_=True)
T(s, Inches(1.15), Inches(4.85), Inches(11.0), Inches(0.5), "青黄不接 → 洗牌红利", size=17, color=WHITE, bold=True, font=F_HEAD)
T(s, Inches(1.15), Inches(5.4), Inches(11.0), Inches(1.0),
  "安利在中国30多年，老一辈六七十岁、打法跟不上新世代；全中国做安利没几个人用AI。\n谁能先跑出“AI新打法”，谁就吃到这波洗牌红利。",
  size=15, color=GOLD_L, font=F_BODY, ls=1.3)
footer(s, 13)

# ==================================================================
# S14 金句总结
# ==================================================================
s = new_slide(NAVY)
rect(s, 0, 0, prs.slide_width, Inches(0.18), GOLD)
T(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8), "金句 · 一页带走", size=30, color=WHITE, bold=True, font=F_HEAD)
rect(s, Inches(0.9), Inches(1.5), Inches(1.6), Inches(0.05), GOLD)
quotes = [
    "安利的公平，是算账算出来的——市场最大的人，收入不一定最高。",
    "先发后面部门的钱，谁都不被辜负。",
    "助人自助，不是利用人。",
    "一栋5层房子月租7500，要花300万买；一组稳健市场，一年不到就能拿到。",
]
for i, q in enumerate(quotes):
    t = Inches(1.5 + i * 1.42)
    rect(s, Inches(0.9), t, Inches(0.07), Inches(1.15), GOLD)
    T(s, Inches(1.2), t + Inches(0.12), Inches(11.3), Inches(0.9), q,
      size=18, color=WHITE, font=F_KAI, ls=1.2)
T(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
  "AI营养师 · 张医生  |  2026-09-01", size=12, color=GOLD_L, font=F_BODY)

# ==================================================================
# SAVE
# ==================================================================
out_dir = '/Users/mac/Documents/zfc最强大脑/A知识沉淀大脑/01-旧知识沉淀/02-创业知识'
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, '安利合作模式-OPP三大奖金分配制度.pptx')
prs.save(out)
print('PPT saved to:', out)
print('Total slides:', len(prs.slides))
