# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
C = {
    'CREAM': RGBColor(0xFA, 0xF3, 0xE0),
    'CREAM_LIGHT': RGBColor(0xFF, 0xF8, 0xEC),
    'HOT_PINK': RGBColor(0xFF, 0x00, 0x66),
    'BRIGHT_YLW': RGBColor(0xFF, 0xD7, 0x00),
    'CYAN': RGBColor(0x00, 0xBF, 0xFF),
    'BLACK': RGBColor(0x22, 0x22, 0x22),
    'WHITE': RGBColor(0xFF, 0xFF, 0xFF),
    'DARK': RGBColor(0x44, 0x44, 0x44),
    'MED': RGBColor(0x99, 0x99, 0x99),
    'LIGHT_BG': RGBColor(0xF0, 0xE8, 0xD8),
    'GRID': RGBColor(0xE0, 0xD8, 0xC8),
    'BLACK_GRID': RGBColor(0xD8, 0xD0, 0xC0),
    'GREEN': RGBColor(0x66, 0xCC, 0x66),
}
FONT = 'Arial'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

def bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C['CREAM']

def tb(slide, l, t, w, h, txt, sz=18, clr=None, bold=False, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = clr or C['BLACK']
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def add_p(tf, txt, sz=16, clr=None, bold=False, align=PP_ALIGN.LEFT, sb=Pt(4), sa=Pt(2)):
    p = tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = clr or C['BLACK']
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    p.space_before = sb
    p.space_after = sa
    return p

def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or C['WHITE']
    shape.line.color.rgb = border or C['BLACK']
    shape.line.width = bw
    return shape

def tag(slide, l, t, txt, bg_c=None, txt_c=None):
    shape = rect(slide, l, t, Inches(1.8), Inches(0.4), fill=bg_c or C['HOT_PINK'], border=C['BLACK'])
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(11)
    p.font.color.rgb = txt_c or C['WHITE']
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    return shape

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    return s

def s_title(slide, num, title, sub=""):
    tag(slide, Inches(0.4), Inches(0.3), "SLIDE " + str(num), bg_c=C['BLACK'], txt_c=C['BRIGHT_YLW'])
    tb(slide, Inches(0.4), Inches(1.0), Inches(12), Inches(0.8), title, sz=32, bold=True)
    if sub:
        tb(slide, Inches(0.4), Inches(1.7), Inches(12), Inches(0.5), sub, sz=14, clr=C['MED'])

def add_bar_chart(slide, left, top, width, items, bar_color=None):
    max_v = max(i[1] for i in items) if items else 100
    bh = Inches(0.35)
    gap = Inches(0.15)
    y = top
    for label, val in items:
        ratio = val / max_v
        tb(slide, left, y, Inches(1.8), bh, label, sz=11, clr=C['DARK'], bold=True, anchor=MSO_ANCHOR.MIDDLE)
        rect(slide, left + Inches(1.9), y, width - Inches(1.9), bh, fill=C['LIGHT_BG'], border=C['BLACK'], bw=Pt(1))
        bar_w = int((width - Inches(1.9)) * ratio)
        rect(slide, left + Inches(1.9), y, bar_w, bh, fill=bar_color or C['HOT_PINK'], border=C['BLACK'], bw=Pt(1))
        tb(slide, left + Inches(1.9) + bar_w + Pt(4), y, Inches(1.2), bh, f"{val}%", sz=11, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        y += bh + gap

# ========== SLIDE 01: COVER ==========
s = new_slide()
# Title
tb(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5), '糖尿病逆转', sz=64, bold=True)
rect(s, Inches(1.5), Inches(3.0), Inches(4), Inches(0.08), fill=C['HOT_PINK'], border=C['BLACK'])
tb(s, Inches(1.5), Inches(3.3), Inches(10), Inches(0.8), 'DIABETES REVERSAL', sz=28, clr=C['HOT_PINK'], bold=True)
tb(s, Inches(1.5), Inches(4.2), Inches(8), Inches(0.6), '科学认知与解决方案 - 循证医学 - 生活方式医学 - 营养干预', sz=16, clr=C['DARK'])
rect(s, Inches(9), Inches(5.5), Inches(4), Inches(2), fill=C['BRIGHT_YLW'], border=C['BLACK'])
tb(s, Inches(9.3), Inches(5.8), Inches(3.5), Inches(1.5), '张大侠\n2026', sz=20, bold=True)
# decor
rect(s, Inches(0), Inches(0), Inches(0.3), Inches(3.5), fill=C['HOT_PINK'], border=C['BLACK'])
for i in range(5):
    rect(s, Inches(12.5), Inches(1.0 + i*0.45), Inches(0.35), Inches(0.35),
         fill=C['CYAN'] if i%2==0 else C['HOT_PINK'], border=C['BLACK'])

# ========== SLIDE 02: TOC ==========
s = new_slide()
s_title(s, '02', 'NAVIGATION', '课程目录')
parts = [
    ('01', '流行病学', '糖尿病发病全景', C['HOT_PINK'], ['中国数据：11.9%患病率', '35.2%糖尿病前期', '慢病四大骑士', '发病机制总览']),
    ('02', '诊断与症状', '识别与评估', C['CYAN'], ['三多一少典型症状', '并发症预警信号', '诊断金标准', 'C肽与胰岛素抵抗']),
    ('03', '解决方案', '逆转路线图', C['BRIGHT_YLW'], ['86%逆转率科学依据', '六大支柱干预', '三位一体：控养提', '案例见证']),
]
for i, (num, title, sub, color, items) in enumerate(parts):
    x = Inches(0.8 + i*4.2)
    tb(s, x, Inches(2.5), Inches(1.5), Inches(1.0), num, sz=60, clr=color, bold=True)
    tb(s, x+Inches(0.8), Inches(2.7), Inches(3), Inches(0.5), title, sz=28, bold=True)
    tb(s, x+Inches(0.8), Inches(3.2), Inches(3), Inches(0.4), sub, sz=13, clr=C['MED'])
    rect(s, x+Inches(0.8), Inches(3.7), Inches(2.5), Inches(0.04), fill=color, border=C['BLACK'])
    for j, item in enumerate(items):
        tb(s, x+Inches(0.8), Inches(4.0+j*0.45), Inches(3.2), Inches(0.4), '> ' + item, sz=12, clr=C['DARK'])

# ========== SLIDE 03: China Data ==========
s = new_slide()
s_title(s, '03', 'DIABETES IN CHINA', '中国成人糖尿病流行病学')
tb(s, Inches(1), Inches(2.5), Inches(5), Inches(1.5), '11.9%', sz=72, clr=C['HOT_PINK'], bold=True)
tb(s, Inches(1), Inches(3.8), Inches(5), Inches(0.5), '中国成人糖尿病患病率', sz=18, bold=True)
tb(s, Inches(1), Inches(4.3), Inches(5), Inches(0.5), '每10人中就有1位糖尿病患者', sz=14, clr=C['MED'])
tb(s, Inches(1), Inches(5.0), Inches(5), Inches(1.5), '35.2%', sz=72, clr=C['CYAN'], bold=True)
tb(s, Inches(1), Inches(6.3), Inches(5), Inches(0.5), '糖尿病前期检出率 - 超1/3为后备军', sz=18, bold=True)
# Info block
block = rect(s, Inches(7), Inches(2.2), Inches(5.5), Inches(4.8))
tf = block.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '关键认知'; p.font.size = Pt(22); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
for t in ['', '控糖不是痛苦的节食与禁忌，', '而是回归身体平衡的科学生活方式。', '', '糖尿病不再是终身的单向列车，', '而是可以科学折返的通道。', '', '来源：中国成人糖尿病流行病学调查', 'DiRECT Trial - The Lancet', '2025 ACLM 指南']:
    add_p(tf, t, sz=14, clr=C['DARK'] if t.startswith('来源') else C['BLACK'], bold='来源' in t)

# ========== SLIDE 04: Four Knights ==========
s = new_slide()
s_title(s, '04', 'CHRONIC DISEASE LANDSCAPE', '慢病四大骑士 - 共享细胞级病理窗口')
diseases = [('心血管病', C['HOT_PINK'], Inches(0.8)), ('糖尿病\n代谢综合征', C['CYAN'], Inches(3.8)),
            ('癌症', C['BRIGHT_YLW'], Inches(6.8)), ('阿尔茨海默', C['HOT_PINK'], Inches(9.8))]
for name, color, x in diseases:
    block = rect(s, x, Inches(2.5), Inches(2.5), Inches(1.8), fill=color)
    tf = block.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(18); p.font.color.rgb = C['WHITE']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except: pass
bottom = rect(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0))
tf = bottom.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '底层共性病理'; p.font.size = Pt(16); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=8)
add_p(tf, '> 线粒体功能障碍 - 细胞发电厂生锈断电', sz=14)
add_p(tf, '> 胰岛素抵抗 - 细胞门禁系统失灵', sz=14)
add_p(tf, '> 慢性低度炎症 - 全身隐性火灾', sz=14)
add_p(tf, '', sz=8)
add_p(tf, 'Reference: Peter Attia Outlive - Robbins Basic Pathology', sz=11, clr=C['MED'])

# ========== SLIDE 05: Pathophysiology ==========
s = new_slide()
s_title(s, '05', 'PATHOPHYSIOLOGY OVERVIEW', '从细胞到宏观的发病链条 - 4步演化')
steps = [
    ('STEP 01', '坏油泛滥-细胞膜僵硬', C['HOT_PINK'], 'Omega-6:Omega-3比例恶化至20:1\n细胞膜弹性丧失-信号接收障碍'),
    ('STEP 02', '胰岛素受体生锈-抵抗', C['CYAN'], '胰岛素狂刷卡-门禁系统被糊死\n糖进不去细胞-血液高糖'),
    ('STEP 03', '代偿性高胰岛素血症', C['BRIGHT_YLW'], '胰腺拼命分泌更多胰岛素\nbeta细胞超负荷-逐渐累瘫'),
    ('STEP 04', '内质网应激-细胞凋亡', C['HOT_PINK'], '细胞内部垃圾堆积-启动自杀程序\n胰岛功能衰竭-不可逆糖尿病'),
]
for i, (label, title, color, desc) in enumerate(steps):
    y = Inches(2.3 + i*1.25)
    if i < 3:
        tb(s, Inches(6.5), y+Inches(0.1), Inches(0.5), Inches(0.5), '↓', sz=24, bold=True, align=PP_ALIGN.CENTER)
    t = rect(s, Inches(0.8), y, Inches(1.8), Inches(0.9), fill=color)
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(13); p.font.color.rgb = C['WHITE']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except: pass
    tb(s, Inches(3.0), y, Inches(3.5), Inches(0.5), title, sz=16, bold=True)
    tb(s, Inches(7.5), y+Inches(0.05), Inches(5), Inches(0.8), desc, sz=12, clr=C['DARK'])

# ========== SLIDE 06: Symptoms ==========
s = new_slide()
s_title(s, '06', 'SYMPTOMS CHECKLIST', '典型症状 - 三多一少')
symptoms = [('多饮', '口渴难止，饮水量大增', '\U0001F4A7', C['HOT_PINK']),
            ('多尿', '夜尿频繁，尿量增多', '\U0001F6BD', C['CYAN']),
            ('多食', '易饿，食量增加', '\U0001F37D', C['BRIGHT_YLW']),
            ('体重下降', '吃得多反而瘦', '⚖', C['HOT_PINK'])]
for i, (name, desc, icon, color) in enumerate(symptoms):
    x = Inches(0.8 + i*3.1)
    rect(s, x, Inches(2.5), Inches(2.7), Inches(2.2))
    ib = rect(s, x+Inches(0.85), Inches(2.7), Inches(1.0), Inches(1.0), fill=color)
    tf = ib.text_frame; p = tf.paragraphs[0]; p.text = icon; p.font.size = Pt(28); p.alignment = PP_ALIGN.CENTER
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except: pass
    tb(s, x+Inches(0.1), Inches(3.8), Inches(2.5), Inches(0.4), name, sz=20, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.1), Inches(4.2), Inches(2.5), Inches(0.4), desc, sz=12, clr=C['DARK'], align=PP_ALIGN.CENTER)
warn = rect(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2), fill=RGBColor(0xFF, 0xF0, 0xE0), border=C['HOT_PINK'])
tf = warn.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '! 很多早期患者无明显症状，体检才发现血糖已超标'; p.font.size = Pt(16); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except: pass

# ========== SLIDE 07: Complications ==========
s = new_slide()
s_title(s, '07', 'COMPLICATION WARNING', '并发症预警信号 - 早发现早干预')
comps = [('视力模糊', '视网膜病变', C['HOT_PINK']), ('四肢麻木', '周围神经病变', C['CYAN']),
         ('伤口难愈', '微循环障碍', C['BRIGHT_YLW']), ('皮肤瘙痒', '高糖刺激神经末梢', C['HOT_PINK']),
         ('泡沫尿', '肾小球损伤信号', C['CYAN'])]
for i, (symp, diag, color) in enumerate(comps):
    y = Inches(2.5 + i*0.85)
    rect(s, Inches(0.8), y, Inches(0.3), Inches(0.6), fill=color, border=C['BLACK'])
    tb(s, Inches(1.4), y, Inches(3), Inches(0.6), '  ' + symp, sz=18, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(5), y, Inches(4), Inches(0.6), '--  ' + diag, sz=14, clr=C['MED'], anchor=MSO_ANCHOR.MIDDLE)
key = rect(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), fill=C['BLACK'])
tf = key.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '糖尿病是全身血管病'; p.font.size = Pt(20); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
add_p(tf, '并发症源于长期高糖对血管的腐蚀 - 控制血糖就是保护血管', sz=14, clr=C['WHITE'], align=PP_ALIGN.CENTER)

# ========== SLIDE 08: Diagnosis ==========
s = new_slide()
s_title(s, '08', 'DIAGNOSTIC CRITERIA', '糖尿病诊断金标准')
criteria = [
    ('空腹血糖', [('>= 7.0 mmol/L', '糖尿病', C['HOT_PINK']), ('6.1-7.0 mmol/L', '糖尿病前期', C['BRIGHT_YLW']), ('3.9-6.1 mmol/L', '正常', C['GREEN'])]),
    ('餐后2h血糖', [('>= 11.1 mmol/L', '糖尿病', C['HOT_PINK']), ('7.8-11.1 mmol/L', '糖尿病前期', C['BRIGHT_YLW']), ('< 7.8 mmol/L', '正常', C['GREEN'])]),
    ('糖化血红蛋白 HbA1c', [('>= 6.5%', '诊断标准', C['HOT_PINK']), ('< 6.5%', '逆转目标', C['GREEN'])]),
]
for i, (title, rows) in enumerate(criteria):
    x = Inches(0.8 + i*4.1)
    block = rect(s, x, Inches(2.3), Inches(3.7), Inches(4.0))
    tf = block.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.name = FONT
    for val, label, color in rows:
        add_p(tf, '', sz=6)
        add_p(tf, val, sz=20, clr=color, bold=True)
        add_p(tf, label, sz=13, clr=C['MED'])

# ========== SLIDE 09: C-Peptide ==========
s = new_slide()
s_title(s, '09', 'C-PEPTIDE ANALYSIS', 'C肽 - 胰岛beta细胞功能的透视镜')
lb = rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
tf = lb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '空腹C肽（正常：0.10-0.43 ng/mL）'; p.font.size = Pt(15); p.font.bold = True; p.font.name = FONT
for val, desc, color in [('> 0.43', '胰岛素抵抗-代偿性亢进', C['HOT_PINK']), ('0.10-0.43', '功能正常-逆转黄金窗口', C['GREEN']), ('< 0.10', '胰岛功能严重受损', C['MED'])]:
    add_p(tf, '', sz=6)
    add_p(tf, val + ' -- ' + desc, sz=14, clr=color, bold=True)
rb = rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5))
tf = rb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '餐后2h C肽 - 需达空腹的3-5倍'; p.font.size = Pt(15); p.font.bold = True; p.font.name = FONT
for val, desc, color in [('< 3倍', '储备功能下降-典型2型糖尿病', C['BRIGHT_YLW']), ('> 5倍', '严重胰岛素抵抗', C['HOT_PINK']), ('几乎不升高', '胰岛功能近乎衰竭', C['MED'])]:
    add_p(tf, '', sz=6)
    add_p(tf, val + ' -- ' + desc, sz=14, clr=color, bold=True)
add_p(tf, '', sz=10)
add_p(tf, '核心：空腹C肽 >= 0.8 + 病程 <= 5年 = 逆转黄金种子', sz=14, clr=C['CYAN'], bold=True)

# ========== SLIDE 10: Insulin Resistance ==========
s = new_slide()
s_title(s, '10', 'INSULIN RESISTANCE', '胰岛素抵抗 - 细胞地铁比喻法')
n = rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.0))
tf = n.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '正常情况'; p.font.size = Pt(16); p.font.color.rgb = C['GREEN']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=6)
add_p(tf, '乘客(血糖)->安保(胰岛素)刷卡->大门(GLUT4)打开->进入车厢(细胞)->发电供能', sz=13)
r = rect(s, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.3))
tf = r.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '胰岛素抵抗'; p.font.size = Pt(16); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=6)
add_p(tf, '细胞膜被Omega-6坏油浸泡->安检机生锈->', sz=13)
add_p(tf, '安保狂刷卡->大门死活不开->', sz=13)
add_p(tf, '乘客堵在站台(高血糖)->车厢内已超载(糖化升高)', sz=13)
sol = rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.6), fill=C['BLACK'])
tf = sol.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '解决方案'; p.font.size = Pt(18); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT
for item, desc, color in [('1 限流', '控糖饮食-减少乘客进站', C['HOT_PINK']), ('2 清空车厢', '促代谢-降糖化-腾出空位', C['CYAN']), ('3 修复安检机', '营养修复-恢复胰岛功能', C['BRIGHT_YLW'])]:
    add_p(tf, '', sz=4)
    add_p(tf, item, sz=16, clr=color, bold=True)
    add_p(tf, '   ' + desc, sz=13, clr=C['WHITE'])

# ========== SLIDE 11: 86% Evidence ==========
s = new_slide()
s_title(s, '11', 'CLINICAL EVIDENCE', '86%临床缓解率 - The Lancet DiRECT试验')
tb(s, Inches(1), Inches(2.3), Inches(4), Inches(2), '86%', sz=96, clr=C['HOT_PINK'], bold=True)
tb(s, Inches(1), Inches(4.2), Inches(4), Inches(0.5), '早期2型糖尿病临床缓解率', sz=18, bold=True)
ev = rect(s, Inches(5.5), Inches(2.3), Inches(7), Inches(4.5))
tf = ev.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'THE LANCET - DiRECT Trial'; p.font.size = Pt(18); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=8)
for t in ['历时7年-多中心随机对照临床试验', '', '2025 ACLM指南重磅共识：', '早期2型糖尿病临床缓解率最高达86%', '', '这不是民间偏方', '不是神医广告', '是严谨的生活方式医学成果', '', '糖尿病不再是终身的单向列车', '而是可以折返的通道']:
    is_imp = '共识' in t or '86%' in t or '成果' in t or '列车' in t or '通道' in t
    add_p(tf, t, sz=14, clr=C['HOT_PINK'] if is_imp else C['BLACK'], bold=is_imp)

# ========== SLIDE 12: Remission Definition ==========
s = new_slide()
s_title(s, '12', 'REMISSION DEFINITION', '临床缓解/逆转的硬性标准')
pre = rect(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.2), fill=C['HOT_PINK'])
tf = pre.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '核心前提：停用所有降糖药及胰岛素'; p.font.size = Pt(20); p.font.color.rgb = C['WHITE']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except: pass
inds = [('糖化血红蛋白 HbA1c', '< 6.5%', C['HOT_PINK']), ('空腹血糖', '< 7.0 mmol/L', C['CYAN']), ('餐后2小时血糖', '< 10.0 mmol/L', C['BRIGHT_YLW'])]
for i, (label, val, color) in enumerate(inds):
    x = Inches(0.8 + i*4.1)
    block = rect(s, x, Inches(4.0), Inches(3.7), Inches(1.6))
    tf = block.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = val; p.font.size = Pt(28); p.font.color.rgb = color; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    add_p(tf, label, sz=14, clr=C['DARK'], align=PP_ALIGN.CENTER)
note = rect(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.8), fill=C['BLACK'])
tf = note.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '需连续维持3个月以上 - 逆转不是盲目停药，是身体恢复自我调节能力'; p.font.size = Pt(14); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except: pass

# ========== SLIDE 13: Golden Candidates ==========
s = new_slide()
s_title(s, '13', 'GOLDEN CANDIDATES', '谁最有希望实现逆转？')
filters = [('黄金窗口期', '确诊2型糖尿病 <= 5年', '⏳', C['HOT_PINK']),
           ('胰岛底线', '空腹C肽 >= 0.8 ng/mL\n（胰岛未坏死，只是休眠）', '\U0001F52C', C['CYAN']),
           ('形体特征', 'BMI >= 24\n或腰围：男>=90cm/女>=85cm', '\U0001F4D0', C['BRIGHT_YLW'])]
for i, (title, desc, icon, color) in enumerate(filters):
    x = Inches(0.8 + i*4.1)
    ib = rect(s, x, Inches(2.5), Inches(3.7), Inches(1.0), fill=color)
    tf = ib.text_frame; p = tf.paragraphs[0]; p.text = icon + ' ' + title; p.font.size = Pt(20); p.font.color.rgb = C['WHITE']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except: pass
    db = rect(s, x, Inches(3.6), Inches(3.7), Inches(1.6))
    tf = db.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(13); p.font.alignment = PP_ALIGN.CENTER; p.font.name = FONT
summ = rect(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2), fill=C['BLACK'])
tf = summ.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '符合以上 - 你就是拥有86%逆转机会的黄金种子选手'; p.font.size = Pt(18); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except: pass

# ========== SLIDE 14: Weight Loss vs Remission ==========
s = new_slide()
s_title(s, '14', 'WEIGHT LOSS - REMISSION', '减重与逆转率的线性法则 - DiRECT Trial')
add_bar_chart(s, Inches(0.8), Inches(2.5), Inches(7), [('< 5 kg', 7), ('5-10 kg', 34), ('10-15 kg', 57), ('>= 15 kg', 86)])
right = rect(s, Inches(8.5), Inches(2.5), Inches(4), Inches(4.0))
tf = right.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '为什么有效？'; p.font.size = Pt(16); p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=8)
add_p(tf, '内脏脂肪包裹胰腺', sz=13, clr=C['DARK'])
add_p(tf, '- beta细胞被闷住', sz=13, clr=C['DARK'])
add_p(tf, '', sz=4)
add_p(tf, '减重-清除内脏脂肪', sz=13, clr=C['CYAN'], bold=True)
add_p(tf, '- beta细胞睡醒', sz=13, clr=C['CYAN'], bold=True)
add_p(tf, '', sz=4)
add_p(tf, '- 重新分泌胰岛素', sz=13, clr=C['DARK'])
add_p(tf, '- 胰岛素抵抗消失', sz=13, clr=C['DARK'])
add_p(tf, '- 血糖回归正常', sz=13, clr=C['GREEN'], bold=True)
tb(s, Inches(0.8), Inches(5.5), Inches(7), Inches(0.5), 'Source: DiRECT Trial - The Lancet', sz=12, clr=C['MED'])

# ========== SLIDE 15: Six Pillars ==========
s = new_slide()
s_title(s, '15', 'SIX PILLARS', '生活方式医学六大支柱 - 缺一不可')
pillars = [('\U0001F37D', '01', '营养精准干预', '核心突破15kg', C['HOT_PINK']), ('\U0001F3C3', '02', '运动双引擎', '肌肉=血糖仓库', C['CYAN']),
           ('\U0001F634', '03', '睡眠管理', '皮质醇控制', C['BRIGHT_YLW']), ('\U0001F9D8', '04', '压力管理', '正念降糖', C['HOT_PINK']),
           ('\U0001F48A', '05', '科学减药', '阶梯式退出', C['CYAN']), ('\U0001F4CA', '06', '血糖监测', '数据导航', C['BRIGHT_YLW'])]
for i, (icon, num, title, sub, color) in enumerate(pillars):
    col, row = i%3, i//3
    x, y = Inches(0.8+col*4.1), Inches(2.3+row*2.5)
    rect(s, x, y, Inches(3.7), Inches(2.1))
    nb = rect(s, x+Inches(0.1), y+Inches(0.1), Inches(0.6), Inches(0.6), fill=color)
    tf = nb.text_frame; p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(14); p.font.color.rgb = C['WHITE']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except: pass
    tb(s, x+Inches(0.9), y+Inches(0.15), Inches(2.5), Inches(0.5), icon+' '+title, sz=18, bold=True)
    tb(s, x+Inches(0.3), y+Inches(0.9), Inches(3.2), Inches(0.4), sub, sz=13, clr=color, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.3), y+Inches(1.4), Inches(3.2), Inches(0.5), '单一维度效果仅30-50%\n六维协同-放大逆转率', sz=11, clr=C['MED'], align=PP_ALIGN.CENTER)

# ========== SLIDE 16: Nutrition ==========
s = new_slide()
s_title(s, '16', 'PILLAR 01: NUTRITION', '营养精准干预 - 进餐革命')
p1 = rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.0))
tf = p1.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '第一阶段：快速破局期（12-20周）'; p.font.size = Pt(16); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=6)
add_p(tf, '极低热量全代营养餐（825-853大卡/天）', sz=14, bold=True)
add_p(tf, '目标：精准控糖控脂，快速清除内脏脂肪', sz=13, clr=C['DARK'])
p2 = rect(s, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.3))
tf = p2.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '第二阶段：长期维持策略'; p.font.size = Pt(16); p.font.color.rgb = C['CYAN']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=6)
add_p(tf, '进餐顺序法则：', sz=14, bold=True)
add_p(tf, '蔬菜 -> 蛋白质 -> 主食', sz=20, clr=C['CYAN'], bold=True)
add_p(tf, '可降低餐后血糖峰值15-20%', sz=13, clr=C['DARK'])
plate = rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.6), fill=C['BLACK'])
tf = plate.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '控糖餐盘法则'; p.font.size = Pt(20); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
add_p(tf, '', sz=10)
add_p(tf, '1/2 蔬菜 - 非淀粉类蔬菜为主', sz=15, clr=C['GREEN'], bold=True, align=PP_ALIGN.CENTER)
add_p(tf, '1/4 蛋白质 - 鱼/去皮禽肉/豆制品', sz=15, clr=C['HOT_PINK'], bold=True, align=PP_ALIGN.CENTER)
add_p(tf, '1/4 粗粮 - 全谷物/杂豆/薯类', sz=15, clr=C['CYAN'], bold=True, align=PP_ALIGN.CENTER)
add_p(tf, '', sz=8)
add_p(tf, '每日主食中粗粮应占至少一半', sz=15, clr=C['BRIGHT_YLW'], bold=True, align=PP_ALIGN.CENTER)
add_p(tf, '两餐之间可吃低GI水果（一拳大小）', sz=14, clr=C['WHITE'], align=PP_ALIGN.CENTER)

# ========== SLIDE 17: Red Lines ==========
s = new_slide()
s_title(s, '17', 'NUTRITION RED LINES', '红线禁区 - 必须远离')
reds = [('添加糖', '含糖饮料、甜点、糖果、果汁', C['HOT_PINK']), ('精制米面', '白米饭、白面条、白馒头、白粥', C['CYAN']),
        ('高GI水果', '西瓜、荔枝、桂圆、熟香蕉', C['BRIGHT_YLW']), ('反式脂肪', '人造奶油、起酥油、反复油炸食品', C['HOT_PINK']),
        ('酒精', '严重抑制肝糖输出', C['CYAN'])]
for i, (title, desc, color) in enumerate(reds):
    y = Inches(2.3+i*0.95)
    block = rect(s, Inches(0.8), y, Inches(5.5), Inches(0.75))
    tf = block.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = 'X  ' + title; p.font.size = Pt(16); p.font.color.rgb = color; p.font.bold = True; p.font.name = FONT
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except: pass
    tb(s, Inches(3.2), y, Inches(3), Inches(0.75), desc, sz=13, clr=C['DARK'], anchor=MSO_ANCHOR.MIDDLE)
rev = rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5))
tf = rev.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '主食革命'; p.font.size = Pt(20); p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=8)
add_p(tf, '推荐的慢糖主食：', sz=14, clr=C['GREEN'], bold=True)
add_p(tf, '全谷物：糙米、燕麦、荞麦、藜麦', sz=13, clr=C['DARK'])
add_p(tf, '杂豆类：红豆、绿豆、鹰嘴豆', sz=13, clr=C['DARK'])
add_p(tf, '薯类：红薯、紫薯、山药', sz=13, clr=C['DARK'])
add_p(tf, '', sz=8)
add_p(tf, '需要警惕的快糖主食：', sz=14, clr=C['HOT_PINK'], bold=True)
add_p(tf, '精制谷物、加工食品、软烂主食', sz=13, clr=C['DARK'])
add_p(tf, '', sz=8)
add_p(tf, '核心原则：粗细搭配，粗粮占一半', sz=15, clr=C['CYAN'], bold=True)

# ========== SLIDE 18: Exercise ==========
s = new_slide()
s_title(s, '18', 'PILLAR 02: EXERCISE', '运动双引擎 - 有氧+抗阻')
cardio = rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.2))
tf = cardio.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '有氧消耗'; p.font.size = Pt(18); p.font.color.rgb = C['CYAN']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=6)
add_p(tf, '每周 >= 150 分钟快走', sz=16, bold=True)
add_p(tf, '最佳策略：三餐后20分钟进行', sz=14, clr=C['DARK'])
add_p(tf, '效果：即时消耗餐后血糖', sz=14, clr=C['CYAN'])
re = rect(s, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.0))
tf = re.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '抗阻训练'; p.font.size = Pt(18); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=6)
add_p(tf, '每周3次-深蹲/靠墙静蹲/弹力带', sz=16, bold=True)
add_p(tf, '效果：增加肌肉量-肌肉=内置血糖仓库', sz=14, clr=C['HOT_PINK'])
add_p(tf, '      根除胰岛素抵抗', sz=14, clr=C['HOT_PINK'])
safe = rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5), fill=C['BLACK'])
tf = safe.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '运动安全须知'; p.font.size = Pt(18); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT
for t, c in [('', None), ('运动前：监测血糖', C['WHITE']), ('  <4.4补碳水 | >16.7暂缓', C['MED']), ('', None),
             ('运动中：常备补给', C['WHITE']), ('  随身带糖果防低血糖', C['MED']), ('', None),
             ('运动后：补充恢复', C['WHITE']), ('  再次监测防延迟性低血糖', C['MED']), ('', None),
             ('原则：循序渐进', C['WHITE']), ('  低强度短时间开始', C['MED'])]:
    if t: add_p(tf, t, sz=13, clr=c, bold=(c==C['WHITE']))
    else: add_p(tf, '', sz=6)

# ========== SLIDE 19: 3-in-1 ==========
s = new_slide()
s_title(s, '19', '3-IN-1 REVERSAL', '三位一体胰岛修复 - 控养提')
steps3 = [
    ('STEP 1', '控-科学管控', '为身体按下暂停键', '通过生活方式管理（饮食、运动）\n平稳血糖-为修复创造窗口期', C['HOT_PINK']),
    ('STEP 2', '养-营养修复', '唤醒细胞的自愈力', '通过精准营养支持\n补充修复原料-唤醒胰岛beta细胞', C['CYAN']),
    ('STEP 3', '提-耐受训练', '重建代谢的弹性空间', '通过糖耐训练循序渐进\n挑战代谢力-重建葡萄糖处理能力', C['BRIGHT_YLW']),
]
for i, (label, title, sub, desc, color) in enumerate(steps3):
    x = Inches(0.8+i*4.1)
    t = rect(s, x, Inches(2.3), Inches(1.5), Inches(0.5), fill=color)
    tf = t.text_frame; p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(12); p.font.color.rgb = C['WHITE']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except: pass
    tb(s, x, Inches(3.0), Inches(3.7), Inches(0.6), title, sz=22, bold=True)
    tb(s, x, Inches(3.6), Inches(3.7), Inches(0.4), sub, sz=14, clr=color, bold=True)
    c = rect(s, x, Inches(4.2), Inches(3.7), Inches(2.3))
    tf = c.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(13); p.font.color.rgb = C['DARK']; p.font.name = FONT
core = rect(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8), fill=C['BLACK'])
tf = core.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '核心理念：不靠药物硬压，唤醒身体本自具足的代谢力'; p.font.size = Pt(16); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except: pass

# ========== SLIDE 20: Green Drink ==========
s = new_slide()
s_title(s, '20', 'PANCREAS REPAIR GREEN DRINK', '胰岛修复绿饮 - 精准营养配方')
ing = rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
tf = ing.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '基础食材'; p.font.size = Pt(16); p.font.bold = True; p.font.name = FONT
for name, amt, effect in [('有机生菜', '50g', '高纤维延缓糖吸收'), ('有机胡萝卜', '50g', 'beta-胡萝卜素滋养胰岛'),
                           ('有机彩椒', '50g', '维C强效抗氧化'), ('牛油果', '半个', '优质脂肪提高胰岛素敏感性')]:
    add_p(tf, '', sz=4)
    add_p(tf, name+' '+amt, sz=14, bold=True)
    add_p(tf, '  --  '+effect, sz=12, clr=C['DARK'])
add_p(tf, '', sz=6)
add_p(tf, '优质油脂：亚麻籽油7.5ml（Omega-3抗炎）', sz=14, clr=C['CYAN'], bold=True)
nut = rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(2.8), fill=C['BLACK'])
tf = nut.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '功能营养素'; p.font.size = Pt(16); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT
for n, a, e, c in [('姜黄维C粉', '10g', '抗炎保护胰岛beta细胞', C['HOT_PINK']), ('肉桂粉', '5g', '激活胰岛素受体', C['CYAN']),
                   ('卵磷脂粉', '10g', '改善脂代谢-缓解抵抗', C['BRIGHT_YLW']), ('多维多矿粉', '10g', '修复必需微量元素', C['WHITE'])]:
    add_p(tf, '', sz=4)
    add_p(tf, n+' '+a, sz=14, clr=c, bold=True)
    add_p(tf, '  --  '+e, sz=12, clr=C['MED'])
meth = rect(s, Inches(7.0), Inches(5.3), Inches(5.5), Inches(1.5))
tf = meth.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '制作：破壁机1分钟 - 立即饮用'; p.font.size = Pt(16); p.font.bold = True; p.font.name = FONT
add_p(tf, '洗净食材-加200ml纯净水-果蔬模式1分钟-即刻饮用', sz=12, clr=C['DARK'])

# ========== SLIDE 21: Sleep & Stress ==========
s = new_slide()
s_title(s, '21', 'PILLAR 03&04: SLEEP & STRESS', '被忽视的逆转关键')
sl = rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
tf = sl.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '睡眠管理'; p.font.size = Pt(20); p.font.color.rgb = C['CYAN']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=8)
add_p(tf, '每日保持7-8小时连续睡眠', sz=16, bold=True)
add_p(tf, '尽量在23:00前入睡', sz=14, clr=C['DARK'])
add_p(tf, '', sz=6)
add_p(tf, '长期熬夜-皮质醇激增', sz=14, clr=C['HOT_PINK'], bold=True)
add_p(tf, '- 胰岛素抵抗增强', sz=14, clr=C['HOT_PINK'])
add_p(tf, '- 内脏脂肪囤积', sz=14, clr=C['HOT_PINK'])
add_p(tf, '', sz=6)
add_p(tf, '改善建议：', sz=14, bold=True)
add_p(tf, ' 固定睡觉和起床时间', sz=13, clr=C['DARK'])
add_p(tf, ' 睡前1小时远离手机', sz=13, clr=C['DARK'])
add_p(tf, ' 暗凉静的卧室环境', sz=13, clr=C['DARK'])
st = rect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5))
tf = st.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '压力管理'; p.font.size = Pt(20); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=8)
add_p(tf, '焦虑压力-皮质醇肾上腺素升高', sz=14, bold=True)
add_p(tf, '- 直接导致血糖升高波动', sz=14, clr=C['HOT_PINK'], bold=True)
add_p(tf, '', sz=6)
add_p(tf, '实用技巧：', sz=14, bold=True)
add_p(tf, ' 每日10分钟正念冥想', sz=13, clr=C['DARK'])
add_p(tf, ' 腹式深呼吸（5分钟快速平复）', sz=13, clr=C['DARK'])
add_p(tf, ' 培养爱好（养花/音乐/书法）', sz=13, clr=C['DARK'])
add_p(tf, ' 社交倾诉-分享感受', sz=13, clr=C['DARK'])
add_p(tf, '', sz=8)
add_p(tf, '这两个支柱常被忽视，', sz=14, bold=True)
add_p(tf, '但它们是逆转能否成功的关键变量', sz=14, bold=True)

# ========== SLIDE 22: Cases ==========
s = new_slide()
s_title(s, '22', 'CASE STUDIES', '真实案例 - 逆转之路')
cases = [('案例一', '37岁女性-病程5年', ['HbA1c: 8.6%-5.6%', '减重15.8kg', '停用全部降糖药和胰岛素', 'C肽恢复正常-2年无反弹'], C['HOT_PINK']),
         ('案例二', '28岁男性-病程5年', ['体重110kg-92kg', 'HbA1c: 8.7%-5.7%', '减重18.0kg', '完全停药-持续缓解3年'], C['CYAN']),
         ('案例三', '更多成效数据', ['TIR从0%提升至100%（1个月）', 'C肽从0.09提升至1.45', '复食高碳水不炸糖', '真正逆转：无需忌碳水'], C['BRIGHT_YLW'])]
for i, (label, sub, items, color) in enumerate(cases):
    x = Inches(0.8+i*4.1)
    block = rect(s, x, Inches(2.3), Inches(3.7), Inches(4.5))
    tf = block.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(18); p.font.color.rgb = color; p.font.bold = True; p.font.name = FONT
    add_p(tf, sub, sz=13, clr=C['DARK'])
    add_p(tf, '', sz=6)
    for item in items:
        add_p(tf, '> '+item, sz=14, bold=True)

# ========== SLIDE 23: Summary ==========
s = new_slide()
s_title(s, '23', 'SUMMARY', '核心信息 - 行动号召')
info = rect(s, Inches(0.8), Inches(2.3), Inches(7), Inches(3.5))
tf = info.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '核心信息'; p.font.size = Pt(22); p.font.color.rgb = C['HOT_PINK']; p.font.bold = True; p.font.name = FONT
add_p(tf, '', sz=6)
add_p(tf, '> 早期2型糖尿病可以逆转', sz=15, bold=True)
add_p(tf, '> 临床缓解率最高达86%', sz=15, clr=C['HOT_PINK'], bold=True)
add_p(tf, '> 不是药物，是生活方式医学', sz=15, clr=C['CYAN'], bold=True)
add_p(tf, '> 控-养-提，三位一体方案', sz=15, clr=C['BRIGHT_YLW'], bold=True)
add_p(tf, '', sz=8)
add_p(tf, '逆转不等于永久根治', sz=15, bold=True)
add_p(tf, '维持健康习惯才能守住成果', sz=13, clr=C['DARK'])
quote = rect(s, Inches(8.5), Inches(2.3), Inches(4), Inches(3.5), fill=C['BLACK'])
tf = quote.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '生活方式医学给予我们的，\n不是一粒神药，\n而是重新掌控自己身体的钥匙'; p.font.size = Pt(18); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except: pass
bot = rect(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8), fill=C['BLACK'])
tf = bot.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '重新掌控自己身体的钥匙'; p.font.size = Pt(18); p.font.color.rgb = C['BRIGHT_YLW']; p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except: pass

# SAVE
out = '/Users/mac/Documents/zfc最强大脑/06-输出（成品）/糖尿病逆转-科学认知与解决方案.pptx'
prs.save(out)
print('PPT saved to:', out)
print('Total slides:', len(prs.slides))
