#!/usr/bin/env python3
"""在「专业版细胞修复营养方案.pptx」总结页前，插入 2 页购买方案 PPT。

风格完全复刻原 deck：
- 米白底 #F5F0E8 / 浅米卡片 #FAF6EF / 纯白卡片 #FFFFFF
- 主色块：黑 #000000 · 青 #00D4E6 · 粉 #FF3A8D · 金 #FFD400
- 正文黑 #1A1A1A / 次级 #444444
- 全程微软雅黑，标题 36pt 加粗，页码 12pt 右下角

插入位置：总结页（「你现在有两个选择」）之前。
输出：专业版细胞修复营养方案-含购买方案.pptx
"""

import re
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 路径 ──
SRC = '/Users/mac/Documents/zfc最强大脑/04-Archive（归档的成品）/01-营养课件/ppt-细胞修复营养方案/专业版细胞修复营养方案.pptx'
OUT = '/Users/mac/Documents/zfc最强大脑/04-Archive（归档的成品）/01-营养课件/ppt-细胞修复营养方案/专业版细胞修复营养方案-含购买方案.pptx'

# ── 色板（与原 deck 完全一致）──
CREAM  = RGBColor(0xF5, 0xF0, 0xE8)   # 背景米白
LIGHT  = RGBColor(0xFA, 0xF6, 0xEF)   # 浅米卡片/隔行
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)   # 正文深色
GRAY   = RGBColor(0x44, 0x44, 0x44)   # 页码/注脚
MUTE   = RGBColor(0x66, 0x66, 0x66)
PINK   = RGBColor(0xFF, 0x3A, 0x8D)   # 粉色警示块
CYAN   = RGBColor(0x00, 0xD4, 0xE6)   # 青色修复块
GOLD   = RGBColor(0xFF, 0xD4, 0x00)   # 金色强调块

FONT = '微软雅黑'

# ── 布局常量（EMU，从原 deck 测量得到）──
CONTENT_X = 457200
CONTENT_W = 11247120
KICKER_Y  = 274320
TITLE_Y   = 731520
FOOTER_Y  = 5640705
FOOTER_H  = 457200
FOOTER_TX_Y = 5686425
PAGENO_X  = 10972800
PAGENO_Y  = 6400800
NOTE_Y    = 6251575

# 表格列 x 坐标（4 列：产品 / 用量 / 零售价 / 折合每月）
COL1_X, COL1_W = 640080, 2514600
COL2_X, COL2_W = 3200400, 1981200
COL3_X, COL3_W = 5227320, 2514600
COL4_X, COL4_W = 7787640, 2514600
HEADER_Y = 1440180
HEADER_H = 457200
ROW_H    = 420048
ROW_GAP  = 420048
ROW0_Y   = 1988820


def set_cjk(run, name=FONT):
    """同时设置拉丁与中文字体为微软雅黑。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        set_cjk(r)
    return tb


def add_rich_text(slide, x, y, w, h, segments, size,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """一段多个 run（用于同一行内强调不同文字）。segments: [(text, bold, color), ...]"""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for text, b, c in segments:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        set_cjk(r)
    return tb


def build_header_and_page(slide, kicker, title, page_str):
    """标题区 + 底部金句条 + 注脚 + 页码（沿用原 deck 固定位置）。"""
    add_text(slide, CONTENT_X, KICKER_Y, 6000000, 457200, kicker, 14, BLACK, bold=True)
    add_text(slide, CONTENT_X, TITLE_Y, 11000000, 914400, title, 36, BLACK, bold=True)
    # 底部金句条（浅米）
    add_rect(slide, CONTENT_X, FOOTER_Y, CONTENT_W, FOOTER_H, LIGHT)
    # 页码
    add_text(slide, PAGENO_X, PAGENO_Y, 1097280, 365760, page_str, 12, GRAY)


def build_slide_a(prs, layout):
    """① 产品价格表：6款产品月花费一览。"""
    slide = prs.slides.add_slide(layout)
    add_bg(slide, CREAM)
    build_header_and_page(slide, '购买方案', '6款产品，月花费一张表算清', '18 / 25')

    # 表头黑条
    add_rect(slide, CONTENT_X, HEADER_Y, CONTENT_W, HEADER_H, BLACK)
    headers = [('产品', COL1_X, COL1_W), ('用量/天', COL2_X, COL2_W),
               ('零售价', COL3_X, COL3_W), ('折合/月', COL4_X, COL4_W)]
    for txt, hx, hw in headers:
        add_text(slide, hx, HEADER_Y + 45720, hw, HEADER_H - 45720, txt, 14, WHITE, bold=True)

    # 产品数据
    products = [
        ('男士活力组合',     '2包',     '868元',      '≈868元'),
        ('多种植物蛋白粉',   '6勺·48g', '560元',      '≈560元'),
        ('针叶樱桃维C',     '3片',     '347元',      '≈347元'),
        ('新款高纯鱼油',     '4粒',     '430元/135粒', '≈383元'),
        ('奶蓟草·组合自带',  '2粒',     '含在组合',   '0元'),
        ('基源欣活',         '1条',     '849元/30条',  '≈849元'),
    ]
    for i, (p, dose, price, month) in enumerate(products):
        y = ROW0_Y + i * ROW_GAP
        fill = LIGHT if i % 2 == 0 else WHITE
        add_rect(slide, CONTENT_X, y, CONTENT_W, ROW_H, fill)
        ty = y + 45720
        add_text(slide, COL1_X, ty, COL1_W, ROW_H - 45720, p, 14, BLACK)
        add_text(slide, COL2_X, ty, COL2_W, ROW_H - 45720, dose, 13, DARK)
        add_text(slide, COL3_X, ty, COL3_W, ROW_H - 45720, price, 13, DARK)
        add_text(slide, COL4_X, ty, COL4_W, ROW_H - 45720, month, 13, DARK)

    # 合计行（黑色，白字加粗）
    ty_total = ROW0_Y + len(products) * ROW_GAP
    add_rect(slide, CONTENT_X, ty_total, CONTENT_W, ROW_H, BLACK)
    tty = ty_total + 45720
    add_text(slide, COL1_X, tty, COL1_W, ROW_H - 45720, '全套修复方案', 14, WHITE, bold=True)
    add_text(slide, COL2_X, tty, COL2_W, ROW_H - 45720, '6款/天', 13, WHITE)
    add_text(slide, COL3_X, tty, COL3_W, ROW_H - 45720, '—', 13, WHITE)
    add_text(slide, COL4_X, tty, COL4_W, ROW_H - 45720, '≈3,007元/月', 14, WHITE, bold=True)

    # 底部金句 + 注脚
    add_text(slide, 731520, FOOTER_TX_Y, 10058400, 365760,
             '全套 ≈3,007元/月 —— 你的细胞，值得这笔投资', 16, BLACK, bold=True)
    add_text(slide, CONTENT_X, NOTE_Y, 9000000, 365760,
             '价格以安利官方/会员实际价为准 · 长客会/会员另有优惠 · 2026年8月', 11, GRAY)
    return slide


def build_slide_b(prs, layout):
    """② 价格总账：每天100元 vs 每天75元 + 省钱通道。"""
    slide = prs.slides.add_slide(layout)
    add_bg(slide, CREAM)
    build_header_and_page(slide, '价格总账', '一天100元，修好全身细胞', '19 / 25')

    # ── 左列：修复账单（青色块）──
    card1_x, card1_y, card1_w, card1_h = CONTENT_X, 1440180, 6400800, 1828800
    add_rect(slide, card1_x, card1_y, card1_w, card1_h, CYAN)
    add_text(slide, card1_x + 274320, card1_y + 114300, 5000000, 457200,
             '修复账单', 18, BLACK, bold=True)
    add_text(slide, card1_x + 274320, card1_y + 480060, 5000000, 914400,
             '≈100元/天', 44, BLACK, bold=True)
    add_text(slide, card1_x + 274320, card1_y + 1400000, 5000000, 457200,
             '3,007元/月 · 全身细胞修复材料', 16, BLACK)

    # ── 左列：伤害账单（粉色块）──
    card2_y = 3359700
    card2_h = 2286000
    add_rect(slide, CONTENT_X, card2_y, card1_w, card2_h, PINK)
    add_text(slide, CONTENT_X + 274320, card2_y + 114300, 5000000, 457200,
             '你每天还在花的「伤害账单」', 18, WHITE, bold=True)
    add_text(slide, CONTENT_X + 274320, card2_y + 457200, 5800000, 1371600,
             '🚬 烟　25元/天 → 750元/月　·　伤血管\n'
             '🧋 奶茶　15元/天 → 450元/月　·　糖+AGEs\n'
             '🍔 外卖　35元/天 → 1,050元/月　·　坏油+美拉德',
             16, WHITE, line_spacing=1.3)
    add_text(slide, CONTENT_X + 274320, card2_y + 1830000, 5800000, 457200,
             '合计 ≈75元/天 · 2,250元/月 —— 每天在「伤害」细胞', 16, WHITE, bold=True)

    # ── 右列：省钱通道（白色卡片 + 金条）──
    card3_x = 7040880
    card3_w = 4663440
    card3_y = 1440180
    card3_h = 4200480
    add_rect(slide, card3_x, card3_y, card3_w, card3_h, WHITE)
    # 顶部金色标题条
    add_rect(slide, card3_x, card3_y, card3_w, 457200, GOLD)
    add_text(slide, card3_x + 91440, card3_y + 45720, 4400000, 365760,
             '省钱通道 —— 实际到手更低', 18, BLACK, bold=True)

    tips = [
        ('注册会员', '　基源欣活 849 → 799'),
        ('满额立减', '　满 1,500 减 380'),
        ('1500长客会', '　全年≈五六折（下页详解）'),
        ('买 6 送 1', '　男士组合 868元/月'),
    ]
    ty = card3_y + 640080
    for tag, desc in tips:
        add_rich_text(slide, card3_x + 91440, ty, 4400000, 457200,
                      [(tag, True, BLACK), (desc, False, DARK)], 15)
        ty += 640080
    add_text(slide, card3_x + 91440, card3_y + 3550000, 4400000, 457200,
             '具体以安利官方 / 顾问方案为准', 13, MUTE)

    # 底部金句 + 注脚
    add_text(slide, 731520, FOOTER_TX_Y, 10058400, 365760,
             '你不是在花钱——你是在把伤害细胞的钱，换成修好细胞的钱', 16, BLACK, bold=True)
    add_text(slide, CONTENT_X, NOTE_Y, 9000000, 365760,
             '价格参考 2026年8月 · 以安利官方/会员实际价格为准', 11, GRAY)
    return slide


def build_slide_c(prs, layout):
    """③ 1500长客会算账：前6个月 vs 全年 回馈率（不签=浪费超一倍）。"""
    slide = prs.slides.add_slide(layout)
    add_bg(slide, CREAM)
    build_header_and_page(slide, '省钱通道', '1500长客会 —— 不签，浪费超一倍', '22 / 25')

    card_w = 5486400
    card_h = 3657600
    cy = 1440180

    # ── 左卡：前6个月（金色）──
    cx1 = 457200
    add_rect(slide, cx1, cy, card_w, card_h, GOLD)
    add_text(slide, cx1 + 274320, cy + 114300, card_w - 548640, 457200,
             '前6个月 · 先算这笔', 18, BLACK, bold=True)
    add_text(slide, cx1 + 274320, cy + 685800, card_w - 548640, 1828800,
             '1-3月积分　1500×1.1×3 ＝ 4,950分\n'
             '4-6月积分　1500×2.2×3 ＝ 9,900分\n'
             '积分兑换　14,850÷5 ＝ 2,970元\n'
             '自用返券　1500×0.9×6%×6 ＝ 486元',
             14, BLACK, line_spacing=1.45)
    add_rect(slide, cx1 + 274320, cy + 2600000, card_w - 548640, 45720, BLACK)
    add_text(slide, cx1 + 274320, cy + 2743200, card_w - 548640, 457200,
             '回馈合计 3,456元', 15, BLACK, bold=True)
    add_text(slide, cx1 + 274320, cy + 3100000, card_w - 548640, 457200,
             '优惠 38% ≈ 六二折', 22, BLACK, bold=True)

    # ── 右卡：全年12个月（青色）──
    cx2 = 6217920
    add_rect(slide, cx2, cy, card_w, card_h, CYAN)
    add_text(slide, cx2 + 274320, cy + 114300, card_w - 548640, 457200,
             '全年12个月 · 更省', 18, BLACK, bold=True)
    add_text(slide, cx2 + 274320, cy + 685800, card_w - 548640, 1828800,
             '总积分　1500×1.1×3＋1500×2.2×9 ＝ 34,650分\n'
             '积分兑换　34,650÷5 ＝ 6,930元\n'
             '自用返券　1500×0.9×6%×12 ＝ 972元\n'
             '──────────',
             14, BLACK, line_spacing=1.45)
    add_rect(slide, cx2 + 274320, cy + 2600000, card_w - 548640, 45720, BLACK)
    add_text(slide, cx2 + 274320, cy + 2743200, card_w - 548640, 457200,
             '回馈合计 7,902元', 15, BLACK, bold=True)
    add_text(slide, cx2 + 274320, cy + 3100000, card_w - 548640, 457200,
             '优惠 44% ≈ 五六折', 22, BLACK, bold=True)

    # ── 底部结论黑条 ──
    add_rect(slide, CONTENT_X, 5151120, CONTENT_W, 365760, BLACK)
    add_text(slide, CONTENT_X, 5217840, CONTENT_W, 274320,
             '同样的产品 —— 不签长客会，资金浪费超一倍', 16, WHITE,
             bold=True, align=PP_ALIGN.CENTER)

    # 底部金句 + 注脚
    add_text(slide, 731520, FOOTER_TX_Y, 10058400, 365760,
             '把积分和返券都拿到手，你的产品才是真·五六折', 16, BLACK, bold=True)
    add_text(slide, CONTENT_X, NOTE_Y, 9000000, 365760,
             '5分=1元 · 积分/返券规则以安利官方当地政策为准 · 2026年8月', 11, GRAY)
    return slide


def build_slide_d(prs, layout):
    """④ 调理周期建议：先修 3-6 个月（肝脏修复周期=一个调理周期），再转正常调理剂量。"""
    slide = prs.slides.add_slide(layout)
    add_bg(slide, CREAM)
    build_header_and_page(slide, '调理建议', '先修 3-6 个月，再转正常调理', '20 / 25')

    card_w = 5486400
    card_h = 3657600
    cy = 1440180

    # ── 左卡：修复期（青色）──
    cx1 = 457200
    add_rect(slide, cx1, cy, card_w, card_h, CYAN)
    add_text(slide, cx1 + 274320, cy + 114300, card_w - 548640, 457200,
             '修复期 · 第1-6个月', 18, BLACK, bold=True)
    add_text(slide, cx1 + 274320, cy + 685800, card_w - 548640, 1828800,
             '参照：肝脏修复周期 5-6 个月\n'
             '＝ 1 个调理周期\n'
             '剂量：修复剂量（6款全套）\n'
             '做法：按每日时间表足量吃',
             14, BLACK, line_spacing=1.45)
    add_rect(slide, cx1 + 274320, cy + 2600000, card_w - 548640, 45720, BLACK)
    add_text(slide, cx1 + 274320, cy + 3100000, card_w - 548640, 457200,
             '先把缺口填满', 20, BLACK, bold=True)

    # ── 右卡：维持期（金色）──
    cx2 = 6217920
    add_rect(slide, cx2, cy, card_w, card_h, GOLD)
    add_text(slide, cx2 + 274320, cy + 114300, card_w - 548640, 457200,
             '维持期 · 6个月后', 18, BLACK, bold=True)
    add_text(slide, cx2 + 274320, cy + 685800, card_w - 548640, 1828800,
             '剂量：转为正常调理剂量\n'
             '节奏：不追求每天全套\n'
             '目的：修复成果不反弹\n'
             '做法：长期稳定维持',
             14, BLACK, line_spacing=1.45)
    add_rect(slide, cx2 + 274320, cy + 2600000, card_w - 548640, 45720, BLACK)
    add_text(slide, cx2 + 274320, cy + 3100000, card_w - 548640, 457200,
             '守住修好的身体', 20, BLACK, bold=True)

    # ── 底部结论黑条 ──
    add_rect(slide, CONTENT_X, 5151120, CONTENT_W, 365760, BLACK)
    add_text(slide, CONTENT_X, 5217840, CONTENT_W, 274320,
             '先足量修 3-6 个月，再转正常调理 —— 一个周期见分晓', 16, WHITE,
             bold=True, align=PP_ALIGN.CENTER)

    # 底部金句 + 注脚
    add_text(slide, 731520, FOOTER_TX_Y, 10058400, 365760,
             '给细胞 3-6 个月，它会还你一个新状态', 16, BLACK, bold=True)
    add_text(slide, CONTENT_X, NOTE_Y, 9000000, 365760,
             '调理周期建议 2026年8月 · 具体以个体情况/顾问建议为准', 11, GRAY)
    return slide


def build_slide_e(prs, layout):
    """⑤ 怎么吃双轨：修复期（按表足量6款）vs 维持期（简化4件套）。"""
    slide = prs.slides.add_slide(layout)
    add_bg(slide, CREAM)
    build_header_and_page(slide, '执行方案', '两个阶段，怎么吃？', '21 / 25')

    cy = 1440180
    card_h = 3657600

    # ── 左卡：修复期（青色，宽）──
    cx1, cw1 = 457200, 6400800
    add_rect(slide, cx1, cy, cw1, card_h, CYAN)
    add_text(slide, cx1 + 274320, cy + 114300, cw1 - 548640, 457200,
             '修复期 · 第1-6个月', 18, BLACK, bold=True)
    add_text(slide, cx1 + 274320, cy + 640080, cw1 - 548640, 457200,
             '修复剂量 · 6款全套', 14, BLACK)
    add_text(slide, cx1 + 274320, cy + 1143000, cw1 - 548640, 2000000,
             '🌅 08:00 早餐　组合1包 + 蛋白粉2勺 + 维C1片 + 鱼油2粒\n'
             '🕛 12:00 午餐　蛋白粉2勺 + 鱼油1粒\n'
             '🕒 16:00 加餐　蛋白粉1勺\n'
             '🌙 18:00 晚餐　组合1包 + 蛋白粉1勺 + 维C2片 + 鱼油1粒\n'
             '🌌 21:30 睡前　基源欣活1条',
             13, BLACK, line_spacing=1.4)
    add_rect(slide, cx1 + 274320, cy + 2900000, cw1 - 548640, 45720, BLACK)
    add_text(slide, cx1 + 274320, cy + 3000000, cw1 - 548640, 500000,
             '🔑 奶蓟草早晚各1粒 · 鱼油随餐更吸收', 13, BLACK, bold=True)

    # ── 右卡：维持期（金色，窄）──
    cx2, cw2 = 7040880, 4663440
    add_rect(slide, cx2, cy, cw2, card_h, GOLD)
    add_text(slide, cx2 + 274320, cy + 114300, cw2 - 548640, 457200,
             '维持期 · 6个月后', 18, BLACK, bold=True)
    add_text(slide, cx2 + 274320, cy + 640080, cw2 - 548640, 457200,
             '正常调理剂量 · 简化4件套', 14, BLACK)
    add_text(slide, cx2 + 274320, cy + 1143000, cw2 - 548640, 2000000,
             '每天一次 · 早餐随餐\n\n'
             '男士组合　1包\n'
             '蛋白粉　3勺\n'
             '维C　1片\n'
             '鱼油　2粒',
             15, BLACK, line_spacing=1.35)
    add_rect(slide, cx2 + 274320, cy + 2900000, cw2 - 548640, 45720, BLACK)
    add_text(slide, cx2 + 274320, cy + 3000000, cw2 - 548640, 500000,
             '简化坚持 · 修复成果不反弹', 13, BLACK, bold=True)

    # ── 底部结论黑条 ──
    add_rect(slide, CONTENT_X, 5151120, CONTENT_W, 365760, BLACK)
    add_text(slide, CONTENT_X, 5217840, CONTENT_W, 274320,
             '修复期足量，维持期简化 —— 两个阶段，两种吃法', 16, WHITE,
             bold=True, align=PP_ALIGN.CENTER)

    # 底部金句 + 注脚
    add_text(slide, 731520, FOOTER_TX_Y, 10058400, 365760,
             '吃对量，才修得到位 —— 别用维持期的量，干修复期的活', 16, BLACK, bold=True)
    add_text(slide, CONTENT_X, NOTE_Y, 9000000, 365760,
             '执行建议 2026年8月 · 具体以顾问建议/个体情况为准', 11, GRAY)
    return slide


def build_slide_f(prs, layout):
    """⑥ 安利奖金比例表 3%→21%：收入比例 + 6%券 + 长客会35% = 综合优惠最高62%。"""
    slide = prs.slides.add_slide(layout)
    add_bg(slide, CREAM)
    build_header_and_page(slide, '额外优惠', '安利奖金比例表 —— 3%到21%，另一重优惠', '24 / 25')

    # 表头黑条（3 列：净营业额 / 收入比例 / 综合优惠；6%券已在长客会页计入，不再重复）
    add_rect(slide, CONTENT_X, HEADER_Y, CONTENT_W, HEADER_H, BLACK)
    headers = [('净营业额(元)', 640080, 3200400), ('收入比例', 3840480, 2514600),
               ('综合优惠(含长客会)', 6355080, 3657600)]
    for txt, hx, hw in headers:
        add_text(slide, hx, HEADER_Y + 45720, hw, HEADER_H - 45720, txt, 14, WHITE, bold=True)

    # 奖金比例数据（档位紧凑排，8 行；综合优惠=长客会12个月44%(含6%券)+收入比例）
    row_h = 384048
    row_gap = 384048
    rows = [
        ('1,250',    '0%',  '≈44%'),
        ('2,500',    '3%',  '≈47%'),
        ('7,500',    '6%',  '≈50%'),
        ('12,500',   '9%',  '≈53%'),
        ('22,500',  '12%',  '≈56%'),
        ('50,000',  '15%',  '≈59%'),
        ('87,500',  '18%',  '≈62%'),
        ('125,000+','21%',  '≈65%'),
    ]
    for i, (bv, rate, total) in enumerate(rows):
        y = ROW0_Y + i * row_gap
        if i == len(rows) - 1:
            add_rect(slide, CONTENT_X, y, CONTENT_W, row_h, GOLD)   # 达标行金色高亮
        else:
            add_rect(slide, CONTENT_X, y, CONTENT_W, row_h,
                     LIGHT if i % 2 == 0 else WHITE)
        ty = y + 45720
        add_text(slide, 640080, ty, 3200400, row_h - 45720, bv, 14, BLACK)
        add_text(slide, 3840480, ty, 2514600, row_h - 45720, rate, 13, DARK)
        add_text(slide, 6355080, ty, 3657600, row_h - 45720, total, 14, BLACK, bold=True)

    # 底部结论黑条
    add_rect(slide, CONTENT_X, 5151120, CONTENT_W, 365760, BLACK)
    add_text(slide, CONTENT_X, 5217840, CONTENT_W, 274320,
             '长客会12个月 + 奖金收入 —— 用产品最高享 ≈65% 优惠', 16, WHITE,
             bold=True, align=PP_ALIGN.CENTER)

    # 金句 + 注脚
    add_text(slide, 731520, FOOTER_TX_Y, 10058400, 365760,
             '别人买产品花钱，你买产品攒优惠 —— 这就是安利的奖金制度', 16, BLACK, bold=True)
    add_text(slide, CONTENT_X, NOTE_Y, 9000000, 365760,
             '达标=月净营业额12.5万 · 综合优惠=长客会12个月44%(含6%券)+收入比例 · 2026年8月', 11, GRAY)
    return slide


def reorder_before_summary(prs):
    """把最后六个 sldId（新页 A-F）移动到总结页之前，顺序 A→B→D→E→C→F。"""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    id_summary = ids[18]      # 第19页 = 总结
    id_a = ids[19]            # 新页A（购买方案）
    id_b = ids[20]            # 新页B（价格总账）
    id_c = ids[21]            # 新页C（长客会算账）
    id_d = ids[22]            # 新页D（调理周期建议）
    id_e = ids[23]            # 新页E（怎么吃双轨）
    id_f = ids[24]            # 新页F（安利奖金比例表）
    for el in (id_summary, id_a, id_b, id_c, id_d, id_e, id_f):
        sldIdLst.remove(el)
    id_before = ids[17]       # 第18页 = 验收
    id_before.addnext(id_a)
    id_a.addnext(id_b)
    id_b.addnext(id_d)
    id_d.addnext(id_e)
    id_e.addnext(id_c)
    id_c.addnext(id_summary)
    id_summary.addnext(id_f)   # 奖金表压轴，放总结之后


def fix_page_numbers(prs):
    """把全库页码分母 18 → 25（仅匹配 'N / 18' 样式）。"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    m = re.fullmatch(r'\s*(\d+)\s*/\s*18\s*', run.text)
                    if m:
                        run.text = f'{m.group(1)} / 25'


def fix_summary_pageno(prs):
    """总结页页码设为 23 / 25（它现在是第24位，奖金表压轴）。

    只针对包含「你现在有两个选择」的总结页，避免误伤其他页。
    """
    for slide in prs.slides:
        has_marker = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = ''.join(r.text for r in para.runs)
                    if '你现在有两个选择' in t:
                        has_marker = True
        if not has_marker:
            continue
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        m = re.fullmatch(r'\s*\d+\s*/\s*25\s*', run.text)
                        if m:
                            run.text = '23 / 25'


def main():
    prs = Presentation(SRC)
    blank = prs.slide_masters[0].slide_layouts[6]  # 无标题/正文占位符
    build_slide_a(prs, blank)
    build_slide_b(prs, blank)
    build_slide_c(prs, blank)
    build_slide_d(prs, blank)
    build_slide_e(prs, blank)
    build_slide_f(prs, blank)
    reorder_before_summary(prs)
    fix_page_numbers(prs)
    fix_summary_pageno(prs)
    prs.save(OUT)
    print('已生成:', OUT)
    print('总页数:', len(prs.slides._sldIdLst))


if __name__ == '__main__':
    main()
