#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI健康 · 超级大脑 — 团队培训PPT
风格：Yellow × Black Editorial（黄底黑字 · 大号现代衬线 · 波普手绘贴纸）
灵感：05-Skills/notebooklm-slide-prompts  →  Yellow × Black Editorial
生成：python-pptx，全文本可编辑
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------------- palette ----------------
Y  = 'FFD500'   # 黄底
YS = 'FFF3A6'   # 浅黄面
YD = 'F0BF00'   # 深黄
BK = '0A0A0A'   # 黑
RD = 'FF3333'   # 红
BL = '0047FF'   # 蓝
WH = 'FFFFFF'   # 白

# ---------------- fonts ----------------
SERIF_CN = 'Songti SC'   # 中文衬线（宋体）
SERIF_EN = 'Didot'       # 英文衬线大标题
SANS_CN  = 'PingFang SC' # 中文正文（黑体）
HAND     = 'Marker Felt' # 手写贴纸
CURSIVE  = 'Snell Roundhand'  # 手写花体

W, H = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def set_font(run, latin=SERIF_EN, ea=SERIF_CN, size=18, bold=True, color=BK, italic=False, spacing=None):
    f = run.font
    f.name = latin
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = RGBColor.from_string(color)
    rPr = run._r.get_or_add_rPr()
    ea_el = rPr.find(qn('a:ea'))
    if ea_el is None:
        ea_el = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea_el)
    ea_el.set('typeface', ea)
    if spacing is not None:
        rPr.set('spc', str(int(spacing * 100)))

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill=None, line=None, lw=2.0,
         shape=MSO_SHAPE.RECTANGLE, rot=0):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if rot:
        sp.rotation = rot
    if fill and fill.upper() != 'NONE':
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def bg(slide, color=Y):
    return rect(slide, 0, 0, W, H, fill=color)

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    return tf

def para(tf, first=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, line_spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    return p

def run(p, txt, latin=SERIF_EN, ea=SERIF_CN, size=18, bold=True, color=BK, italic=False, spacing=None):
    r = p.add_run()
    r.text = txt
    set_font(r, latin=latin, ea=ea, size=size, bold=bold, color=color, italic=italic, spacing=spacing)
    return r

def est_w(text, size):
    w = 0.0
    for ch in text:
        w += size if ord(ch) > 0x2E80 else size * 0.56
    return w / 72 + 0.26

def sticker(slide, x, y, text, fill=WH, rot=-2, size=14, bold=True,
            latin=HAND, ea=SERIF_CN, color=BK, lw=2.0, boxw=None):
    w = boxw if boxw else est_w(text, size)
    h = size / 72 * 1.7 + 0.16
    if fill and fill.upper() != 'NONE':
        rect(slide, x + 0.07, y + 0.07, w, h, fill=BK, rot=rot)
    sp = rect(slide, x, y, w, h, fill=fill, line=BK, lw=lw, rot=rot)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, Inches(0.02))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, text, latin=latin, ea=ea, size=size, bold=bold, color=color)
    return sp

def kicker(slide, text, x=0.45, y=0.5, color=BK, size=15):
    rect(slide, x, y + 0.04, 0.22, 0.22, fill=RD, shape=MSO_SHAPE.OVAL)
    tf = tb(slide, x + 0.36, y, 10, 0.4)
    p = para(tf, first=True)
    run(p, text, ea=SERIF_CN, latin=SERIF_EN, size=size, bold=True, color=color, spacing=0.04)

def title_mixed(slide, parts, x=0.45, y=1.35, w=12.4, size=44, line_below=True, line_w=1.15):
    tf = tb(slide, x, y, w, size / 72 * 1.5 + 0.3)
    p = para(tf, first=True)
    for t, c in parts:
        run(p, t, ea=SERIF_CN, latin=SERIF_EN, size=size, bold=True, color=c)
    if line_below:
        rect(slide, x + 0.02, y + size / 72 * 1.18 + 0.08, line_w, 0.07, fill=BK)
    return tf

def title(slide, text, x=0.45, y=1.35, w=12.4, size=44, color=BK, line_below=True):
    return title_mixed(slide, [(text, color)], x=x, y=y, w=w, size=size, line_below=line_below)

def chrome(slide, n):
    tf = tb(slide, 0.45, 7.08, 8, 0.32)
    p = para(tf, first=True)
    run(p, 'AI营养师 · 张医生', ea=SANS_CN, latin=SANS_CN, size=9, bold=True, color=BK)
    run(p, '   SUPER BRAIN · 团队培训', latin=CURSIVE, ea=SANS_CN, size=12, bold=False, color=BK)
    c = rect(slide, 12.35, 6.78, 0.6, 0.6, fill=BK, shape=MSO_SHAPE.OVAL)
    tf2 = c.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf2, m, 0)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run(p2, str(n), latin=SERIF_EN, size=14, bold=True, color=WH)

def halftone(slide, x, y, cols=12, rows=4, gap=0.22, size=0.05, color=YD):
    for r in range(rows):
        for c in range(cols):
            rect(slide, x + c * gap, y + r * gap, size, size, fill=color, shape=MSO_SHAPE.OVAL)

def card(slide, x, y, w, h, fill=YS, title='', body='', rot=0, accent=BK,
         title_size=20, body_size=14):
    if fill and fill.upper() != 'NONE':
        rect(slide, x + 0.06, y + 0.06, w, h, fill=BK, rot=rot)
    sp = rect(slide, x, y, w, h, fill=fill, line=BK, lw=2.0, rot=rot)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for m in ('margin_left', 'margin_right'):
        setattr(tf, m, Inches(0.16))
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run(p, title, ea=SERIF_CN, latin=SERIF_EN, size=title_size, bold=True, color=accent)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    p2.line_spacing = 1.12
    run(p2, body, ea=SANS_CN, latin=SANS_CN, size=body_size, bold=False, color=BK)
    return sp

def big_stat(slide, x, y, num, unit, label, num_size=110, color=BK):
    tf = tb(slide, x, y, 5.2, 2.8)
    p = para(tf, first=True, align=PP_ALIGN.LEFT)
    run(p, num, latin=SERIF_EN, size=num_size, bold=True, italic=True, color=color)
    if unit:
        run(p, ' ' + unit, latin=SERIF_EN, ea=SERIF_CN, size=int(num_size * 0.42), bold=True, color=color)
    p2 = para(tf, align=PP_ALIGN.LEFT, space_before=8)
    run(p2, label, ea=SERIF_CN, latin=SANS_CN, size=17, bold=True, color=BK)

# =========================================================
# Slide 1 — 封面
# =========================================================
s = add_slide()
bg(s, Y)
# 右下网点
halftone(s, 0.6, 6.0, cols=10, rows=3, gap=0.24)
# 右上红圈 + 黑斜带
rect(s, 10.9, -1.8, 3.6, 3.6, fill='NONE', line=RD, lw=5, shape=MSO_SHAPE.OVAL)
rect(s, 11.7, -2.0, 2.3, 12, fill=BK, rot=28)
# kicker
kicker(s, 'AI 营养师 · 张医生', size=16)
# English display
tf = tb(s, 0.68, 1.55, 9, 0.6)
p = para(tf, first=True)
run(p, 'HEALTH × AI', latin=SERIF_EN, size=26, bold=True, italic=True, color=RD, spacing=0.06)
run(p, '  SUPER BRAIN', latin=SERIF_EN, size=26, bold=True, italic=False, color=BK)
# 主标题
tf = tb(s, 0.62, 2.05, 11.5, 3.6)
p = para(tf, first=True)
run(p, 'AI健康', ea=SERIF_CN, latin=SERIF_EN, size=96, bold=True, color=BK)
p2 = para(tf)
run(p2, '超级大脑', ea=SERIF_CN, latin=SERIF_EN, size=96, bold=True, color=BK)
rect(s, 0.68, 5.05, 1.7, 0.1, fill=RD)
# 副标题
tf = tb(s, 0.7, 5.35, 10, 0.6)
p = para(tf, first=True)
run(p, '人人都是 AI营养师 · 1000+ 团队，一起冲 FC', ea=SERIF_CN, latin=SANS_CN, size=22, bold=True, color=BK)
# 贴纸
sticker(s, 0.95, 6.15, '医生 × AI × 系统', fill=WH, rot=-3, size=14)
sticker(s, 3.35, 6.45, '2026 团队培训', fill=RD, rot=2, size=14, color=WH)
sticker(s, 6.1, 6.2, '风口 · 正在发生', fill=WH, rot=-1.5, size=14)

# =========================================================
# Slide 2 — 我是谁
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 11.6, 0.5, cols=6, rows=3, gap=0.22)
kicker(s, '01 · 自我介绍')
title_mixed(s, [('我是 ', BK), ('张医生', RD)], size=50)
tf = tb(s, 0.47, 2.28, 9, 0.5)
p = para(tf, first=True)
run(p, '一个 AI 营养师 — 把三样东西，合而为一', ea=SERIF_CN, latin=SANS_CN, size=18, bold=True, color=BK)
# 三张卡片
cw, ch = 3.7, 2.5
card(s, 0.5, 3.0, cw, ch, fill=WH, title='医生专业', body='南昌大学医学院 · 三甲医院麻醉医生\n懂身体，更懂健康的底层逻辑', accent=BK)
sticker(s, 4.5, 3.85, '×', fill='NONE', rot=0, size=30, latin=SERIF_EN, color=RD, boxw=0.7)
card(s, 5.05, 3.0, cw, ch, fill=YS, title='AI 科技', body='用 AI 把知识变成生产力\n一个人，就是一支团队', accent=RD)
sticker(s, 9.05, 3.85, '×', fill='NONE', rot=0, size=30, latin=SERIF_EN, color=RD, boxw=0.7)
card(s, 9.6, 3.0, cw, ch, fill=WH, title='安利系统', body='6 年全职安利\n带团队从 0 到 1000+', accent=BL)
# 公式
tf = tb(s, 0.6, 5.85, 12, 0.7)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '医生专业 ', ea=SERIF_CN, latin=SERIF_EN, size=20, bold=True, color=BK)
run(p, '×', latin=SERIF_EN, size=24, bold=True, color=RD)
run(p, ' AI科技 ', ea=SERIF_CN, latin=SERIF_EN, size=20, bold=True, color=BK)
run(p, '×', latin=SERIF_EN, size=24, bold=True, color=RD)
run(p, ' 安利系统', ea=SERIF_CN, latin=SERIF_EN, size=20, bold=True, color=BK)
run(p, '  =   AI营养师', ea=SERIF_CN, latin=SERIF_EN, size=22, bold=True, color=RD)
sticker(s, 11.15, 6.4, '三合为一', fill=RD, rot=2, size=14, color=WH)
chrome(s, 2)

# =========================================================
# Slide 3 — 我的历程
# =========================================================
s = add_slide()
bg(s, Y)
kicker(s, '01 · 自我介绍')
title(s, '从放牛娃 到 AI营养师', size=46)
sticker(s, 9.3, 1.2, '6年没上班 · 全靠安利', fill=RD, rot=2, size=14, color=WH)
sticker(s, 9.35, 6.15, '每一步 都算数', fill=WH, rot=-2, size=13)
# 时间线
steps = [
    ('农村娃', '干农活 · 放牛'),
    ('省重点高中', '奥赛班'),
    ('南昌大学医学院', '211 · 学医'),
    ('三甲医院', '麻醉医生'),
    ('26岁创业', '开奶茶店 · 亏30多万'),
    ('全职安利', '6年 · 从0开始'),
    ('现在', 'AI营养师 · 超级大脑'),
]
line_x = 1.05
rect(s, line_x, 1.75, 0.045, 4.85, fill=BK)
y0 = 1.78
for i, (t, d) in enumerate(steps):
    yy = y0 + i * 0.68
    big = (i == len(steps) - 1)
    fill = RD if big else BK
    c = rect(s, line_x - 0.16, yy, 0.34, 0.34, fill=fill, shape=MSO_SHAPE.OVAL)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, str(i + 1), latin=SERIF_EN, size=12, bold=True, color=WH)
    tf = tb(s, 1.6, yy - 0.06, 4.4, 0.45)
    p = para(tf, first=True)
    run(p, t, ea=SERIF_CN, latin=SERIF_EN, size=18, bold=True, color=BK)
    if d:
        tf = tb(s, 1.6, yy + 0.3, 4.6, 0.35)
        p = para(tf, first=True)
        run(p, d, ea=SANS_CN, latin=SANS_CN, size=12, bold=False, color=BK)
# 右侧大字
tf = tb(s, 7.0, 2.6, 5.6, 3.2)
p = para(tf, first=True, align=PP_ALIGN.RIGHT)
run(p, '6年', latin=SERIF_EN, size=120, bold=True, italic=True, color=RD)
p2 = para(tf, align=PP_ALIGN.RIGHT)
run(p2, '全职 · 全靠安利', ea=SERIF_CN, latin=SANS_CN, size=20, bold=True, color=BK)
p3 = para(tf, align=PP_ALIGN.RIGHT, space_before=10)
run(p3, '以前踩过的坑，现在变成带团队的底气', ea=SANS_CN, latin=SANS_CN, size=14, bold=False, color=BK)
chrome(s, 3)

# =========================================================
# Slide 4 — 我的目标
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 11.8, 0.55, cols=5, rows=3, gap=0.22)
kicker(s, '02 · 我的目标')
# FC 大圈
c = rect(s, 1.0, 1.7, 3.1, 3.1, fill='NONE', line=BK, lw=4, shape=MSO_SHAPE.OVAL)
rect(s, 1.28, 2.18, 2.5, 2.5, fill=RD, shape=MSO_SHAPE.OVAL)
tf = tb(s, 1.28, 2.32, 2.5, 1.6, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, 'FC', latin=SERIF_EN, size=92, bold=True, color=WH)
# 右侧标题
title_mixed(s, [('皇冠大使', BK)], x=4.7, y=1.9, w=8.6, size=56, line_w=1.3)
tf = tb(s, 4.7, 3.1, 7.8, 1.2)
p = para(tf, first=True)
run(p, '一个人走得快，一群人走得远', ea=SERIF_CN, latin=SERIF_EN, size=22, bold=True, color=BK)
p2 = para(tf, space_before=8)
run(p2, '用 AI营养师，带 1000+ 团队，一起冲 FC', ea=SERIF_CN, latin=SANS_CN, size=17, bold=True, color=BK)
sticker(s, 9.7, 4.7, '我的北极星', fill=WH, rot=-3, size=15)
# 底部三个锚点
anchors = [('目标', '皇冠大使 FC'), ('路径', 'AI营养师'), ('团队', '1000+ 一起走')]
ax = 0.6
for t, d in anchors:
    card(s, ax, 5.6, 3.55, 1.1, fill=YS, title=t, body=d, title_size=15, body_size=16)
    ax += 3.95
chrome(s, 4)

# =========================================================
# Slide 5 — AI初步成果
# =========================================================
s = add_slide()
bg(s, Y)
kicker(s, '03 · AI初步成果')
title_mixed(s, [('AI ', BK), ('真的能出结果', RD)], size=46)
sticker(s, 11.35, 1.15, '风口已验证', fill=WH, rot=2, size=13)
# 两大数据
big_stat(s, 0.75, 2.5, '18,000', '粉', '1个月 · 自媒体涨粉', num_size=108)
big_stat(s, 6.4, 2.5, '4', '条', '10万+ 播放 · 连续爆款', num_size=108, color=RD)
# 上升箭头
arr = rect(s, 9.9, 2.6, 1.9, 1.9, fill=BK, shape=MSO_SHAPE.UP_ARROW)
# 说明卡
card(s, 0.75, 5.55, 11.9, 1.15, fill=WH,
     title='不是我多厉害，是 AI 把普通人的能力放大了',
     body='自媒体 × AI视频 · 一条视频 = 一个完整团队的制作能力',
     title_size=19, body_size=14)
chrome(s, 5)

# =========================================================
# Slide 6 — 风口导航
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 12.0, 6.1, cols=5, rows=2, gap=0.22)
kicker(s, '04 · AI视频风口')
title_mixed(s, [('风口 ', BK),('已经来了', RD)], size=46)
tf = tb(s, 0.47, 2.28, 9, 0.5)
p = para(tf, first=True)
run(p, '5 个板块，一次讲透', ea=SERIF_CN, latin=SANS_CN, size=18, bold=True, color=BK)
# 5 行列表
rows = [
    ('01', '本地知识库', '我的第二大脑 — 把读过的书、听过的课全部沉淀'),
    ('02', '资料库 · 文件检索', '一句话，AI 帮你把资料找出来'),
    ('03', '知识的门槛', '被 AI 砍掉了 — 人人可调用顶尖知识'),
    ('04', 'Vibe Coding', '2025 年首次提出 — 普通人的 AI 入口'),
    ('05', '普通人的机遇', '现在，就是上车的最好时机'),
]
yy = 2.75
for num, t, d in rows:
    rect(s, 0.62, yy, 0.72, 0.72, fill=BK)
    tf = tb(s, 0.62, yy, 0.72, 0.72, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, num, latin=SERIF_EN, size=22, bold=True, color=WH)
    tf = tb(s, 1.55, yy - 0.04, 10.8, 0.75)
    p = para(tf, first=True)
    run(p, t, ea=SERIF_CN, latin=SERIF_EN, size=21, bold=True, color=BK)
    run(p, '   ' + d, ea=SANS_CN, latin=SANS_CN, size=14, bold=False, color=BK)
    yy += 0.86
sticker(s, 10.85, 5.6, '不是预测 · 是正在发生', fill=RD, rot=2, size=13, color=WH)
chrome(s, 6)

# =========================================================
# Slide 7 — 本地知识库
# =========================================================
s = add_slide()
bg(s, Y)
kicker(s, '04 · AI视频风口  ①')
title(s, '本地知识库', size=52)
tf = tb(s, 0.47, 2.35, 9, 0.6)
p = para(tf, first=True)
run(p, '我的第二大脑', ea=SERIF_CN, latin=SANS_CN, size=22, bold=True, color=RD)
# 三个内容卡
items = [('读过的书', '知识不丢'), ('听过的课', '精华都在'), ('讲过的内容', '一键复用')]
ax = 0.6
for t, d in items:
    card(s, ax, 3.1, 3.6, 1.5, fill=WH, title=t, body=d, title_size=22, body_size=14)
    ax += 3.9
# 主体句
tf = tb(s, 0.62, 5.0, 11.8, 1.1)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '把一切知识，沉淀成 AI 可检索的资产', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=BK)
p2 = para(tf, align=PP_ALIGN.CENTER, space_before=8)
run(p2, '知识永不丢失 · 随用随取 · 越用越值钱', ea=SANS_CN, latin=SANS_CN, size=16, bold=False, color=BK)
sticker(s, 11.2, 6.15, '第二大脑', fill=RD, rot=-2, size=15, color=WH)
chrome(s, 7)

# =========================================================
# Slide 8 — 资料库 · 文件检索
# =========================================================
s = add_slide()
bg(s, Y)
kicker(s, '04 · AI视频风口  ②')
title(s, 'AI 帮你找资料', size=52)
tf = tb(s, 0.47, 2.35, 9, 0.6)
p = para(tf, first=True)
run(p, '文件检索 · 秒级响应', ea=SERIF_CN, latin=SANS_CN, size=22, bold=True, color=RD)
# 对话气泡
bub = rect(s, 0.7, 3.0, 6.2, 0.95, fill=WH, line=BK, lw=2)
tf = bub.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.2)
p = tf.paragraphs[0]
run(p, '你：  ', ea=SERIF_CN, latin=SERIF_EN, size=16, bold=True, color=BK)
run(p, '“帮我把那份资料找出来”', ea=SANS_CN, latin=SANS_CN, size=16, bold=False, color=BK)
rect(s, 0.62, 3.78, 0.42, 0.42, fill=BK, shape=MSO_SHAPE.OVAL)
bub2 = rect(s, 0.7, 4.35, 6.2, 0.95, fill=BK, line=BK, lw=2)
tf2 = bub2.text_frame
tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
tf2.margin_left = Inches(0.2)
p2 = tf2.paragraphs[0]
run(p2, 'AI：  ', latin=SERIF_EN, ea=SERIF_CN, size=16, bold=True, color=RD)
run(p2, '“找到了，在这里 ↓”', ea=SANS_CN, latin=SANS_CN, size=16, bold=False, color=WH)
# 右侧要点
tf = tb(s, 7.6, 3.0, 5.0, 2.6)
p = para(tf, first=True)
run(p, '不用再翻文件夹', ea=SERIF_CN, latin=SERIF_EN, size=22, bold=True, color=BK)
p2 = para(tf, space_before=12)
run(p2, '一句话，资料自己出来', ea=SERIF_CN, latin=SERIF_EN, size=22, bold=True, color=BK)
p3 = para(tf, space_before=12)
run(p3, '本地资料库 × AI = 你的私人研究员', ea=SANS_CN, latin=SANS_CN, size=15, bold=False, color=BK)
sticker(s, 8.0, 6.1, '秒级响应', fill=WH, rot=2, size=14)
sticker(s, 10.5, 6.45, '告别大海捞针', fill=RD, rot=-2, size=13, color=WH)
chrome(s, 8)

# =========================================================
# Slide 9 — 知识门槛被砍掉
# =========================================================
s = add_slide()
bg(s, Y)
kicker(s, '04 · AI视频风口  ③')
title_mixed(s, [('知识的门槛 ', BK), ('被 AI 砍掉了', RD)], size=46, line_w=1.5)
# 前后对比
card(s, 0.7, 2.6, 5.5, 2.4, fill=WH, title='过去',
     body='知识被垄断\n专业遥不可及\n想学，找不到入口', title_size=26, body_size=16)
tf = tb(s, 6.35, 3.35, 1.0, 1.0, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '→', latin=SERIF_EN, size=40, bold=True, color=RD)
card(s, 7.15, 2.6, 5.5, 2.4, fill=BK, title='现在',
     body='AI 把知识的门槛砍掉了\n人人可调用顶尖知识\n入口，已经打开', title_size=26, body_size=16, accent=RD)
# 修正黑色卡片正文颜色
# 底部金句
tf = tb(s, 0.7, 5.5, 11.9, 0.9)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '你缺的从来不是智商，是入口', ea=SERIF_CN, latin=SERIF_EN, size=30, bold=True, color=BK)
p2 = para(tf, align=PP_ALIGN.CENTER, space_before=6)
run(p2, '而 AI，把入口免费打开了', ea=SANS_CN, latin=SANS_CN, size=16, bold=False, color=BK)
sticker(s, 10.9, 6.3, '入口 已打开', fill=RD, rot=-2, size=14, color=WH)
chrome(s, 9)

# =========================================================
# Slide 10 — Vibe Coding
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 0.5, 6.2, cols=8, rows=2, gap=0.22)
kicker(s, '04 · AI视频风口  ④')
title_mixed(s, [('Vibe Coding', RD)], x=0.45, y=1.4, size=54, line_w=1.4)
tf = tb(s, 0.47, 2.5, 9, 0.6)
p = para(tf, first=True)
run(p, '2025 年 首次提出', latin=SERIF_EN, ea=SERIF_CN, size=24, bold=True, italic=True, color=BK)
run(p, '  普通人的 AI 入口', ea=SERIF_CN, latin=SANS_CN, size=18, bold=True, color=BK)
# 大 2025
tf = tb(s, 8.2, 1.0, 4.4, 2.2)
p = para(tf, first=True, align=PP_ALIGN.RIGHT)
run(p, '2025', latin=SERIF_EN, size=110, bold=True, italic=True, color=RD)
# 解释
card(s, 0.7, 3.3, 7.6, 1.3, fill=WH, title='用「聊天」就能让 AI 帮你干活',
     body='不需要会编程 · 说出你的需求，AI 帮你实现', title_size=20, body_size=15)
card(s, 0.7, 4.85, 7.6, 1.3, fill=YS, title='这是“普通人用上 AI”的标志概念',
     body='Vibe = 状态 · Coding = 创造 —— 跟着感觉，说出想法', title_size=20, body_size=15)
# 右侧要点
tf = tb(s, 8.9, 3.4, 3.9, 2.4)
p = para(tf, first=True)
run(p, '一句话', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=BK)
p2 = para(tf, space_before=8)
run(p2, '描述需求', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=BK)
p3 = para(tf, space_before=8)
run(p3, 'AI 帮你干活', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=RD)
sticker(s, 10.9, 6.15, '会说话 · 就能用', fill=WH, rot=-3, size=14)
chrome(s, 10)

# =========================================================
# Slide 11 — 普通人的 AI 机遇
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 0.5, 0.55, cols=8, rows=2, gap=0.22)
kicker(s, '04 · AI视频风口  ⑤')
title_mixed(s, [('每个普通人 ', BK), ('都能上车', RD)], size=50)
tf = tb(s, 0.47, 2.4, 11.9, 0.6)
p = para(tf, first=True)
run(p, '不是懂技术的人赢，是会用 AI 的人赢', ea=SERIF_CN, latin=SERIF_EN, size=22, bold=True, color=BK)
# 两句话
tf = tb(s, 0.7, 3.3, 11.8, 0.7)
p = para(tf, first=True)
run(p, '①  AI 让普通人，站在同一起跑线', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=BK)
tf = tb(s, 0.7, 4.2, 11.8, 0.7)
p = para(tf, first=True)
run(p, '②  你的经验 + AI 的能力 = 你的新杠杆', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=BK)
# 大标语
rect(s, 0.7, 5.35, 11.9, 1.0, fill=BK)
tf = tb(s, 0.7, 5.35, 11.9, 1.0, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '现在，就是最好的上车时机', ea=SERIF_CN, latin=SERIF_EN, size=28, bold=True, color=WH)
sticker(s, 11.3, 1.15, '弯道超车', fill=RD, rot=2, size=14, color=WH)
sticker(s, 0.7, 6.5, '你，准备好了吗？', fill=WH, rot=-1.5, size=13)
chrome(s, 11)

# =========================================================
# Slide 12 — 未来 · AI+自媒体
# =========================================================
s = add_slide()
bg(s, Y)
kicker(s, '05 · 未来已来')
title_mixed(s, [('AI ', BK), ('×', RD), (' 自媒体', BK)], size=52)
tf = tb(s, 0.47, 2.4, 9, 0.6)
p = para(tf, first=True)
run(p, '人人都是自媒体', ea=SERIF_CN, latin=SANS_CN, size=22, bold=True, color=RD)
# 三个要点
items = [('内容生产效率', '×10', '一人顶一个团队'),
         ('创作门槛', '归零', '拿起手机就能拍'),
         ('分发渠道', '全通', '短视频打穿所有信息')]
ax = 0.6
for t, num, d in items:
    card(s, ax, 3.0, 3.75, 2.2, fill=WH, title=t, body=d, title_size=18, body_size=14)
    tf = tb(s, ax, 3.45, 3.75, 1.0, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, num, latin=SERIF_EN, size=44, bold=True, italic=True, color=RD)
    ax += 4.0
tf = tb(s, 0.7, 5.6, 11.9, 0.7)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '一条视频 = 一个人 + AI = 一支完整团队', ea=SERIF_CN, latin=SERIF_EN, size=24, bold=True, color=BK)
sticker(s, 10.85, 6.3, '你的手机 = 你的电视台', fill=RD, rot=-2, size=13, color=WH)
chrome(s, 12)

# =========================================================
# Slide 13 — 人人都是百万网红
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 12.0, 0.6, cols=5, rows=2, gap=0.22)
kicker(s, '05 · 未来已来')
title_mixed(s, [('人人都是 ', BK), ('百万网红', RD)], size=50, line_w=1.5)
tf = tb(s, 0.47, 2.5, 11.9, 0.6)
p = para(tf, first=True)
run(p, '这个时代，已经把所有路都铺好了', ea=SERIF_CN, latin=SERIF_EN, size=20, bold=True, color=BK)
# 两块
card(s, 0.7, 3.2, 5.75, 1.6, fill=WH, title='专业信息',
     body='已被短视频完全打通\n任何专业知识，都能被讲给普通人听', title_size=20, body_size=15)
card(s, 6.85, 3.2, 5.75, 1.6, fill=WH, title='运营方法论',
     body='短视频运营方法也已全部打通\n有方法可抄，有路径可走', title_size=20, body_size=15)
# 大标语
rect(s, 0.7, 5.25, 11.9, 1.05, fill=BK)
tf = tb(s, 0.7, 5.25, 11.9, 1.05, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '普通人做网红的时代，一定会到来', ea=SERIF_CN, latin=SERIF_EN, size=28, bold=True, color=WH)
sticker(s, 0.95, 6.5, '而你，要提前上车', fill=WH, rot=-2, size=13)
sticker(s, 10.9, 6.5, '时代 来了', fill=RD, rot=2, size=14, color=WH)
chrome(s, 13)

# =========================================================
# Slide 14 — 我的超级大脑
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 0.5, 0.55, cols=7, rows=2, gap=0.22)
kicker(s, '06 · 超级大脑')
title_mixed(s, [('我的 ', BK), ('超级大脑', RD)], size=50)
tf = tb(s, 0.47, 2.4, 9, 0.6)
p = para(tf, first=True)
run(p, '知识库 × AI × 自媒体 = 超级个体', ea=SERIF_CN, latin=SERIF_EN, size=20, bold=True, color=BK)
# 三个组件
parts = [('① 本地知识库', '资产', '知识永远是你的'),
         ('② AI 工具', '生产力', '一个人的团队'),
         ('③ 短视频', '放大器', '让世界看见你')]
ax = 0.6
for t, tag, d in parts:
    card(s, ax, 3.1, 3.75, 2.0, fill=WH, title=t, body=d, title_size=20, body_size=14)
    sticker(s, ax + 1.0, 4.25, tag, fill=RD, rot=-2, size=13, color=WH)
    ax += 4.0
tf = tb(s, 0.7, 5.6, 11.9, 0.7)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '给每个普通人，配一个超级大脑', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=BK)
sticker(s, 11.2, 6.3, '人人都能拥有', fill=WH, rot=2, size=13)
chrome(s, 14)

# =========================================================
# Slide 15 — 结尾
# =========================================================
s = add_slide()
bg(s, Y)
halftone(s, 11.8, 0.55, cols=5, rows=3, gap=0.22)
# 大 FC
rect(s, 10.6, 3.4, 3.3, 3.3, fill='NONE', line=BK, lw=4, shape=MSO_SHAPE.OVAL)
rect(s, 10.85, 3.85, 2.6, 2.6, fill=RD, shape=MSO_SHAPE.OVAL)
tf = tb(s, 10.85, 4.0, 2.6, 1.7, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, 'FC', latin=SERIF_EN, size=84, bold=True, color=WH)
# 主文案
kicker(s, '07 · 一起出发')
tf = tb(s, 0.62, 1.7, 10, 1.6)
p = para(tf, first=True)
run(p, '成为', ea=SERIF_CN, latin=SERIF_EN, size=42, bold=True, color=BK)
run(p, ' AI营养师', ea=SERIF_CN, latin=SERIF_EN, size=42, bold=True, color=RD)
p2 = para(tf, space_before=10)
run(p2, '让每个普通人都能拥有超级大脑', ea=SERIF_CN, latin=SERIF_EN, size=26, bold=True, color=BK)
tf = tb(s, 0.7, 4.2, 10, 0.8)
p = para(tf, first=True)
run(p, 'AI健康 · 超级大脑 · 我们一起去', ea=SERIF_CN, latin=SANS_CN, size=20, bold=True, color=BK)
sticker(s, 0.95, 5.55, '加入我们', fill=WH, rot=-3, size=15)
sticker(s, 2.75, 5.85, '一起冲 FC', fill=RD, rot=2, size=15, color=WH)
tf = tb(s, 0.7, 6.7, 8, 0.5)
p = para(tf, first=True)
run(p, '谢谢 · 未来见', ea=SERIF_CN, latin=CURSIVE, size=20, bold=True, color=BK)
chrome(s, 15)

# ---------------- save ----------------
out = '/Users/mac/Documents/zfc最强大脑/04-Archive（归档的成品）/培训PPT-AI营养师超级大脑/AI健康培训-超级大脑-黄黑编辑风.pptx'
prs.save(out)
print('SAVED:', out)
print('slides:', len(prs.slides.__iter__.__self__._sldIdLst))
