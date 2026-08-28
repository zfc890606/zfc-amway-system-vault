# -*- coding: utf-8 -*-
"""张飞梦想清单 · 杂志风重制 —— 内容不变，风格：明朝体+手写混搭·黄底黑字·时尚杂志感"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

YELLOW = RGBColor(0xFF, 0xD6, 0x00)
YELLOW_D = RGBColor(0xF0, 0xC2, 0x00)
BLACK   = RGBColor(0x11, 0x11, 0x11)
RED     = RGBColor(0xD7, 0x00, 0x2F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x55, 0x55, 0x55)
F_SERIF = '宋体'
F_KAI   = '楷体'
F_SANS  = '黑体'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
N_TOTAL = 9

def _cjk(run, font_cn):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_cn)

def _style_run(r, text, size, color, bold, font):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    _cjk(r, font)

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = YELLOW
    return s

def page_chrome(s, idx):
    pbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(prs.slide_width * (idx+1) / N_TOTAL), Inches(0.045))
    pbar.fill.solid(); pbar.fill.fore_color.rgb = RED; pbar.line.fill.background()
    # 四角裁切标记
    for (cx, cy, sx, sy) in [(0.12,0.12,1,1),(12.13,0.12,-1,1),(0.12,7.28,1,-1),(12.13,7.28,-1,-1)]:
        m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(0.16*sx), Inches(0.16))
        m.fill.solid(); m.fill.fore_color.rgb = BLACK; m.line.fill.background()
    tb(s, Inches(11.7), Inches(7.02), Inches(1.3), Inches(0.4),
       f"{idx+1:02d} — {N_TOTAL:02d}", size=11, color=BLACK, font=F_SANS, align=PP_ALIGN.RIGHT)

def tb(slide, l, t, w, h, text, size=18, color=BLACK, bold=False, font=F_SERIF,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, rotation=0, ls=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = anchor
    except Exception: pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    r = p.add_run()
    _style_run(r, text, size, color, bold, font)
    if rotation: box.rotation = rotation
    return box

def tb_rich(slide, l, t, w, h, parts, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, rotation=0, ls=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try: tf.vertical_anchor = anchor
    except Exception: pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    for (t, sz, c, b, f) in parts:
        r = p.add_run()
        _style_run(r, t, sz, c, b, f)
    if rotation: box.rotation = rotation
    return box

def rect(slide, l, t, w, h, fill=WHITE, line=BLACK, lw=1.5, rotation=0, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if round_:
        try: shp.adjustments[0] = 0.12
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    if rotation: shp.rotation = rotation
    return shp

def oline(slide, l, t, w, h, lw=2.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    shp.fill.background()
    shp.line.color.rgb = BLACK
    shp.line.width = Pt(lw)
    return shp

def sticker(slide, l, t, w, h, text, size=13, rotation=-3, fill=WHITE,
            color=BLACK, bold=False, font=F_KAI):
    shp = rect(slide, l, t, w, h, fill=fill, line=BLACK, lw=1.5, rotation=rotation)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _style_run(r, text, size, color, bold, font)
    return shp

def seal(slide, l, t, size, text, rotation=3):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, size, size)
    try: shp.adjustments[0] = 0.18
    except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = RED
    shp.line.color.rgb = BLACK; shp.line.width = Pt(1.25)
    shp.rotation = rotation
    tf = shp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('vert', 'eaVert')
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    size_in = size / 914400.0 if hasattr(size, 'emu') else float(size)
    _style_run(r, text, max(6, int(size_in * 72 / 3.4)), WHITE, True, F_KAI)
    return shp

# ========== S1 封面 ==========
s = new_slide(); page_chrome(s, 0)
oline(s, Inches(9.0), Inches(-2.4), Inches(6.2), Inches(6.2), lw=2.5)
oline(s, Inches(10.3), Inches(-1.1), Inches(6.2), Inches(6.2), lw=1.25)
tb(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4),
   "张飞的梦想清单", size=14, color=BLACK)
sticker(s, Inches(10.5), Inches(0.5), Inches(2.2), Inches(0.5), "2024.10.03",
        size=13, rotation=3, fill=WHITE, bold=True)
seal(s, Inches(11.7), Inches(5.9), Inches(0.85), "梦想", rotation=3)
tb_rich(s, Inches(1.0), Inches(2.3), Inches(11.5), Inches(2.2),
        [("张", 92, BLACK, True, F_SERIF), ("飞", 96, RED, True, F_KAI)])
rect(s, Inches(1.05), Inches(4.5), Inches(2.2), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(1.05), Inches(4.75), Inches(10), Inches(0.6),
   "一个普通医生的梦想清单", size=22, color=BLACK, font=F_KAI)
tb(s, Inches(1.0), Inches(6.15), Inches(9), Inches(0.6),
   "百花丛中过  片叶不沾身", size=20, color=GRAY, font=F_KAI, rotation=-3)

# ========== S2 一个普通的医生（2012），有什么梦想？ ==========
s = new_slide(); page_chrome(s, 1)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "DREAM LIST", size=12, color=GRAY, font=F_SANS)
sticker(s, Inches(10.9), Inches(0.55), Inches(1.9), Inches(0.55), "开 篇",
        size=14, rotation=4, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.8),
        [("一个普通的医生（", 46, BLACK, True, F_SERIF),
         ("2012", 52, RED, True, F_SERIF),
         ("），\n有什么梦想？", 46, BLACK, True, F_SERIF)])
rect(s, Inches(0.95), Inches(4.6), Inches(1.8), Inches(0.05), fill=BLACK, line=None)
seal(s, Inches(11.5), Inches(5.5), Inches(0.85), "初心", rotation=-3)

# ========== S3-8 梦想清单页 ==========
dreams = [
    ("03", "环游中国", "点亮地图"),
    ("04", "四次进藏", "硬座摩旅 · 单车骑行"),
    ("05", "钻石演出", "记录人生"),
    ("06", "一所房子", "面朝大海"),
    ("07", "千次演讲", "绽放价值"),
    ("08", "下乡支教", "涤荡灵魂"),
]
for idx, (num, main, punch) in enumerate(dreams, start=2):
    s = new_slide(); page_chrome(s, idx)
    # 右侧超大背景数字
    tb(s, Inches(9.2), Inches(1.4), Inches(4.2), Inches(4.5), num,
       size=250, color=YELLOW_D, bold=True, align=PP_ALIGN.RIGHT)
    # 左上小标
    tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), f"DREAM {num}", size=12, color=GRAY, font=F_SANS)
    # 左侧数字红标
    tb(s, Inches(0.9), Inches(1.6), Inches(2.2), Inches(1.4), num,
       size=64, color=RED, bold=True)
    rect(s, Inches(0.95), Inches(3.1), Inches(1.6), Inches(0.05), fill=BLACK, line=None)
    # 主文案
    tb(s, Inches(0.9), Inches(3.5), Inches(9), Inches(1.6), main,
       size=54, color=BLACK, bold=True)
    tb(s, Inches(0.92), Inches(5.0), Inches(9), Inches(1.0), punch,
       size=34, color=BLACK, font=F_KAI)
    # 印章
    seal(s, Inches(11.7), Inches(5.75), Inches(0.8), "梦", rotation=-3)

# ========== S9 结尾 ==========
s = new_slide(); page_chrome(s, 8)
tb(s, Inches(0.9), Inches(0.6), Inches(4), Inches(0.4), "THE END · 未完待续", size=12, color=GRAY, font=F_SANS)
sticker(s, Inches(10.7), Inches(0.55), Inches(2.1), Inches(0.55), "这一生",
        size=14, rotation=4, fill=WHITE, bold=True)
tb_rich(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.7),
        [("此生，", 54, BLACK, True, F_SERIF),
         ("做一个有故事的人", 56, RED, True, F_KAI)])
rect(s, Inches(0.95), Inches(3.8), Inches(2.2), Inches(0.05), fill=BLACK, line=None)
tb(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.9),
   "千磨万击还坚劲，任尔东西南北风", size=28, color=BLACK, font=F_KAI)
tb(s, Inches(0.9), Inches(5.6), Inches(8), Inches(0.6),
   "—— 郑燮《竹石》", size=14, color=GRAY, font=F_KAI)
seal(s, Inches(11.6), Inches(5.4), Inches(0.9), "有故事", rotation=3)

# ========== SAVE ==========
out_dir = '/Users/mac/Documents/zfc最强大脑/04-Archive（归档的成品）/02-创业认知'
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, '张飞梦想清单-杂志风.pptx')
prs.save(out)
print('PPT saved to:', out)
print('Total slides:', len(prs.slides))
